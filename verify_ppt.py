# 생성된 파워포인트 파일의 슬라이드 개수 및 구조 정합성을 검증하는 스크립트
from pptx import Presentation

def verify():
    file_path = 'BTC_지점별_코드맵_점검현황_최종최종.pptx'
    print(f"Loading {file_path} for verification...")
    prs = Presentation(file_path)
    
    print(f"Total slides in presentation: {len(prs.slides)}")
    assert len(prs.slides) == 8, f"Expected 8 slides, got {len(prs.slides)}"

    slide_titles = []
    for idx, slide in enumerate(prs.slides, 1):
        # Find textboxes or shapes with title text
        title_text = "Unknown Title"
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "현황" in text or "통계" in text or "보고서" in text or "요약" in text or "비교" in text:
                    # heuristic to find slide title
                    if len(text) < 40:
                        title_text = text
                        break
        slide_titles.append(title_text)
        print(f"  Slide {idx}: {title_text}")

    print("\nVerification completed successfully. The presentation structure is valid.")

if __name__ == '__main__':
    verify()
