# 3바디 7코드 AI코드맵 신뢰성 홍보 PPTX 브로셔 생성 스크립트
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 테마 색상 (신뢰감을 주는 프리미엄 딥 네이비 & 아일랜드 블루 테마)
COLOR_DARK_BG = RGBColor(15, 23, 42)      # 슬레이트 블랙/딥 네이비
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # 은은한 그레이시 화이트
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PRIMARY = RGBColor(27, 54, 93)      # 다크 네이비
COLOR_SECONDARY = RGBColor(59, 130, 246)  # 아일랜드 블루
COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # 차콜 블랙
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 뮤트 그레이
COLOR_ACCENT = RGBColor(245, 158, 11)     # 골드/오렌지 액센트
COLOR_CARD_BORDER = RGBColor(226, 232, 240)
COLOR_LIGHT_BLUE = RGBColor(239, 246, 255)

FONT_NAME = "Malgun Gothic"

def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title, subtitle):
    # 상단 대분류/서브 타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.3))
    p_sub = sub_box.text_frame.paragraphs[0]
    p_sub.text = subtitle.upper()
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(9.5)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.6))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_DARK

def create_card_shape(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER):
    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = bg_color
    card_shape.line.color.rgb = border_color
    card_shape.line.width = Pt(1)
    return card_shape

def add_paragraph_to_tf(tf, text, font_size=11, is_bold=False, font_color=COLOR_TEXT_DARK, space_after=6, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    return p

def format_cell_content(cell, text, font_size=10, is_bold=False, font_color=COLOR_TEXT_DARK, alignment=PP_ALIGN.CENTER):
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color

def style_table_header(table, col_count):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

def generate_presentation():
    # PPT 생성 시작 (Widescreen 16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==================================================
    # SLIDE 1: 표지 슬라이드
    # ==================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_DARK_BG)

    left_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_SECONDARY
    left_bar.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(2.8), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_SECONDARY
    badge.line.fill.background()
    badge.text_frame.text = "통합 웰니스 건강 진단 솔루션"
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = FONT_NAME
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_WHITE

    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "임직원 건강경영 및 대국민 멘탈 웰니스 확산 프로젝트"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(16)
    p0.font.color.rgb = COLOR_SECONDARY
    p0.font.bold = True
    p0.space_after = Pt(12)

    p1 = tf.add_paragraph()
    p1.text = "3바디 7코드 AI 코드맵 신뢰성 브로셔"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "몸 · 마음 · 뇌 통합 건강나이 진단과 7코드 에너지 솔루션의 과학적 타당성 분석"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_WHITE
    p2.font.bold = False
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "공동 연구 개발: 국가공인 브레인트레이너 전문가 단체 · 한국뇌과학연구원 · 글로벌사이버대학교"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(6)

    p4 = tf.add_paragraph()
    p4.text = "제공 기관: 브레인트레이닝센터"
    p4.font.name = FONT_NAME
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_TEXT_MUTED

    # ==================================================
    # SLIDE 2: 왜 몸·마음·뇌의 통합적 건강인가?
    # ==================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_LIGHT_BG)
    add_slide_header(s2, "건강의 새로운 관점: 신체·마음·뇌의 통합적 균형", "Integrated Health Needs")

    card_w2 = Inches(3.6)
    card_h2 = Inches(5.1)

    # 1. 몸 (신체 건강에 치중된 관리의 한계)
    create_card_shape(s2, Inches(0.8), Inches(1.5), card_w2, card_h2)
    tb_b1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b1 = tb_b1.text_frame
    tf_b1.word_wrap = True
    tf_b1.paragraphs[0].text = "💪 신체 위주 관리의 한계"
    tf_b1.paragraphs[0].font.name = FONT_NAME
    tf_b1.paragraphs[0].font.size = Pt(15)
    tf_b1.paragraphs[0].font.bold = True
    tf_b1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b1.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b1, "• 반쪽짜리 건강 관리", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b1, "대부분의 사람들은 건강이라고 하면 주로 헬스, 요가, 식단 등 육체적인 관리에만 전념합니다. 하지만 몸만 가꾼다고 해서 삶의 활력이 저절로 채워지지는 않습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b1, "• 유기적 연관성", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b1, "몸의 긴장은 뇌의 스트레스를 유발하고, 지친 뇌는 다시 신체 소화기계 및 자율신경계 저하로 이어집니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2. 마음 (소홀하기 쉬운 마음의 병)
    create_card_shape(s2, Inches(4.8), Inches(1.5), card_w2, card_h2)
    tb_b2 = s2.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b2 = tb_b2.text_frame
    tf_b2.word_wrap = True
    tf_b2.paragraphs[0].text = "❤️ 방치되는 마음의 피로와 스트레스"
    tf_b2.paragraphs[0].font.name = FONT_NAME
    tf_b2.paragraphs[0].font.size = Pt(15)
    tf_b2.paragraphs[0].font.bold = True
    tf_b2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b2.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b2, "• 일상 스트레스 인지율 25.9% (질병관리청)", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b2, "성인 4명 중 1명은 일상 중 심한 스트레스를 느끼며, 특히 주 경제 활동층인 30대와 40대의 스트레스 인지가 매우 높게 나타납니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b2, "• 직장인 번아웃 경험률 60~70% 돌파", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b2, "직무 및 대인관계 스트레스로 인해 정서적 탈진 상태를 경험한 비율은 직장인 대상 여러 통계 조사에서 심각한 수치로 드러나고 있습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3. 뇌 (우려와 대책의 공백)
    create_card_shape(s2, Inches(8.8), Inches(1.5), card_w2, card_h2)
    tb_b3 = s2.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b3 = tb_b3.text_frame
    tf_b3.word_wrap = True
    tf_b3.paragraphs[0].text = "🧠 뇌 피로와 인지 노화 방치"
    tf_b3.paragraphs[0].font.name = FONT_NAME
    tf_b3.paragraphs[0].font.size = Pt(15)
    tf_b3.paragraphs[0].font.bold = True
    tf_b3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b3.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b3, "• 스마트폰 과의존과 디지털 치매", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b3, "과도한 디지털 기기 사용으로 현대인의 전두엽은 끊임없이 혹사당하고 있습니다. 이는 주의집중력 저하와 단기 기억력 감퇴라는 브레인 피로로 이어집니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b3, "• 예방적 뇌 관리의 공백", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b3, "기존 건강검진은 구조적 이상(종양, 뇌졸중 등)만 잡아낼 뿐, 일상적인 인지 기능 저하나 스트레스 누적으로 인한 뇌 기능 저하는 평가하지 못합니다.", 10.5, False, COLOR_TEXT_MUTED, 6)


    # ==================================================
    # SLIDE 3: 3BODY 통합 4대 나이 진단 시스템
    # ==================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_LIGHT_BG)
    add_slide_header(s3, "3BODY 다차원 측정: 4대 나이 진단 시스템 구성", "4-Age Diagnosis System")

    card_w3 = Inches(2.7)
    card_h3 = Inches(5.1)

    # 1. 몸나이
    create_card_shape(s3, Inches(0.6), Inches(1.5), card_w3, card_h3)
    tb_n1 = s3.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(2.4), Inches(4.7))
    tf_n1 = tb_n1.text_frame
    tf_n1.word_wrap = True
    tf_n1.paragraphs[0].text = "💪 몸나이 (Physical)"
    tf_n1.paragraphs[0].font.name = FONT_NAME
    tf_n1.paragraphs[0].font.size = Pt(14)
    tf_n1.paragraphs[0].font.bold = True
    tf_n1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_n1.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf_n1, "• 핵심 측정 항목", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n1, "- 거북목 각도\n- 척추 측만 정렬\n- 한발 서기 균형 감각\n- 스쿼트 하체 근력\n- 푸시업 상체 근력\n- 전굴 유연성", 10, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_n1, "• 진단 지표", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n1, "신체 비대칭성과 가동범위, 근력/근지구력을 복합 연산하여 신체적 기능 상태를 나이로 정량화합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 2. 뇌나이
    create_card_shape(s3, Inches(3.7), Inches(1.5), card_w3, card_h3)
    tb_n2 = s3.shapes.add_textbox(Inches(3.85), Inches(1.7), Inches(2.4), Inches(4.7))
    tf_n2 = tb_n2.text_frame
    tf_n2.word_wrap = True
    tf_n2.paragraphs[0].text = "🧠 뇌나이 (Cognitive)"
    tf_n2.paragraphs[0].font.name = FONT_NAME
    tf_n2.paragraphs[0].font.size = Pt(14)
    tf_n2.paragraphs[0].font.bold = True
    tf_n2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_n2.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf_n2, "• 핵심 측정 항목", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n2, "- 시각 인지 반응 속도\n- Stroop 충동 억제\n- 다중 항목 작업 기억력\n- 좌우뇌 균형 반응", 10, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_n2, "• 진단 지표", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n2, "인지 처리 속도 및 정보 선택·제어 전두엽 기능을 평가하여 인지 신경망 노화도를 도출합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 3. 마음나이
    create_card_shape(s3, Inches(6.8), Inches(1.5), card_w3, card_h3)
    tb_n3 = s3.shapes.add_textbox(Inches(6.95), Inches(1.7), Inches(2.4), Inches(4.7))
    tf_n3 = tb_n3.text_frame
    tf_n3.word_wrap = True
    tf_n3.paragraphs[0].text = "❤️ 마음나이 (Emotional)"
    tf_n3.paragraphs[0].font.name = FONT_NAME
    tf_n3.paragraphs[0].font.size = Pt(14)
    tf_n3.paragraphs[0].font.bold = True
    tf_n3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_n3.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf_n3, "• 핵심 측정 항목", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n3, "- 7코드 에너지 분포\n- 스트레스 저항성\n- 감정적 탈진도\n- 무의식 불안 분석", 10, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_n3, "• 진단 지표", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n3, "체형 및 인지 데이터의 패턴 분석을 통해 자율신경계 병목 지점을 찾고 정서적 번아웃 나이를 산출합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 4. 얼굴나이
    create_card_shape(s3, Inches(9.9), Inches(1.5), card_w3, card_h3)
    tb_n4 = s3.shapes.add_textbox(Inches(10.05), Inches(1.7), Inches(2.4), Inches(4.7))
    tf_n4 = tb_n4.text_frame
    tf_n4.word_wrap = True
    tf_n4.paragraphs[0].text = "👤 얼굴나이 (Skin & Facial)"
    tf_n4.paragraphs[0].font.name = FONT_NAME
    tf_n4.paragraphs[0].font.size = Pt(14)
    tf_n4.paragraphs[0].font.bold = True
    tf_n4.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_n4.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf_n4, "• 핵심 측정 항목", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n4, "- 안면 비대칭도\n- 피부 주름 분포\n- 탄력도 및 처짐 현상\n- 안면 혈류 건강 상태", 10, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_n4, "• 진단 지표", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_n4, "컴퓨터 비전을 활용해 눈가, 팔자 주름 깊이와 피부 처짐을 분석하여 외적인 생물학적 에이징 수준을 측정합니다.", 10, False, COLOR_TEXT_MUTED, 6)


    # ==================================================
    # SLIDE 4: 4대 나이 진단별 학술적/과학적 근거
    # ==================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, COLOR_LIGHT_BG)
    add_slide_header(s4, "4대 나이 측정 시스템의 학술적 · 과학적 설계 근거", "Scientific Evidence of 4-Ages")

    card_w4 = Inches(5.6)
    card_h4 = Inches(2.4)

    # 1. 몸나이 근거
    create_card_shape(s4, Inches(0.8), Inches(1.5), card_w4, card_h4)
    tb_e1 = s4.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(5.3), Inches(2.2))
    tf_e1 = tb_e1.text_frame
    tf_e1.word_wrap = True
    tf_e1.paragraphs[0].text = "💪 몸나이: 스포츠의학 표준 및 근감소 검증"
    tf_e1.paragraphs[0].font.name = FONT_NAME
    tf_e1.paragraphs[0].font.size = Pt(13)
    tf_e1.paragraphs[0].font.bold = True
    tf_e1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_e1.paragraphs[0].space_after = Pt(8)
    add_paragraph_to_tf(tf_e1, "• ACSM(미국스포츠의학회) 및 국민체력100 데이터베이스 기반", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e1, "연령대별 근력 및 평형성 표준 오차 범위를 모델링하여 근지구력 지표를 객관적으로 역산합니다.", 10, False, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_e1, "• 사코페니아(근감소증) 및 평형성 지수 노화 상관관계", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e1, "눈 감고 외발 서기 능력은 소뇌 및 자율신경계 평형 기능을 반영하는 대표적인 노화 척도입니다.", 10, False, COLOR_TEXT_MUTED, 0)

    # 2. 뇌나이 근거
    create_card_shape(s4, Inches(6.9), Inches(1.5), card_w4, card_h4)
    tb_e2 = s4.shapes.add_textbox(Inches(7.05), Inches(1.6), Inches(5.3), Inches(2.2))
    tf_e2 = tb_e2.text_frame
    tf_e2.word_wrap = True
    tf_e2.paragraphs[0].text = "🧠 뇌나이: 전두엽 실행 기능 및 인지 통제력"
    tf_e2.paragraphs[0].font.name = FONT_NAME
    tf_e2.paragraphs[0].font.size = Pt(13)
    tf_e2.paragraphs[0].font.bold = True
    tf_e2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_e2.paragraphs[0].space_after = Pt(8)
    add_paragraph_to_tf(tf_e2, "• Stroop(스트룹) 과제 및 전두엽 억제 제어 모델", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e2, "시각 자극 간섭 하에 충동성 반응 속도(ms 단위)를 측정하여 전두엽의 실행 기능을 지표화합니다.", 10, False, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_e2, "• 인지 자원 고갈과 백질 신경망 노화성 감소", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e2, "다중 정보 처리를 담당하는 작업 기억력은 전두엽 피질 두께 감소 및 시냅스 가소성 저하와 정비례합니다.", 10, False, COLOR_TEXT_MUTED, 0)

    # 3. 마음나이 근거
    create_card_shape(s4, Inches(0.8), Inches(4.2), card_w4, card_h4)
    tb_e3 = s4.shapes.add_textbox(Inches(0.95), Inches(4.3), Inches(5.3), Inches(2.2))
    tf_e3 = tb_e3.text_frame
    tf_e3.word_wrap = True
    tf_e3.paragraphs[0].text = "❤️ 마음나이: PNI 및 자율신경계 긴장 분석"
    tf_e3.paragraphs[0].font.name = FONT_NAME
    tf_e3.paragraphs[0].font.size = Pt(13)
    tf_e3.paragraphs[0].font.bold = True
    tf_e3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_e3.paragraphs[0].space_after = Pt(8)
    add_paragraph_to_tf(tf_e3, "• 심신신경면역학(PNI)에 근거한 자율신경 활성도 매핑", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e3, "체형 비대칭, 호흡 근육 긴장 상태를 바탕으로 심박변이도(HRV)와 조응하는 정서 상태를 유추합니다.", 10, False, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_e3, "• 만성 스트레스로 인한 HPA 축 과부하 지수화", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e3, "에너지 코드의 방전 수준은 호르몬 불균형과 정신적 탈진(번아웃)의 심도를 나타냅니다.", 10, False, COLOR_TEXT_MUTED, 0)

    # 4. 얼굴나이 근거
    create_card_shape(s4, Inches(6.9), Inches(4.2), card_w4, card_h4)
    tb_e4 = s4.shapes.add_textbox(Inches(7.05), Inches(4.3), Inches(5.3), Inches(2.2))
    tf_e4 = tb_e4.text_frame
    tf_e4.word_wrap = True
    tf_e4.paragraphs[0].text = "👤 얼굴나이: 외적 생물학적 노화 표지자 분석"
    tf_e4.paragraphs[0].font.name = FONT_NAME
    tf_e4.paragraphs[0].font.size = Pt(13)
    tf_e4.paragraphs[0].font.bold = True
    tf_e4.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_e4.paragraphs[0].space_after = Pt(8)
    add_paragraph_to_tf(tf_e4, "• 컴퓨터 비전 및 딥러닝 기반 안면 랜드마크 분석", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e4, "눈가 및 입가 팔자 주름 깊이, 피부 결 처짐 벡터를 기반으로 안면 노화 나이를 추정합니다.", 10, False, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_e4, "• 안면 대칭도와 순환 저하 징후 추적", 10.5, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_e4, "안면 신경망 긴장 및 미세 대칭 붕괴는 내적 만성 피로와 뇌 피로의 중요한 피부 생리적 발현 지표입니다.", 10, False, COLOR_TEXT_MUTED, 0)


    # ==================================================
    # SLIDE 5: 통합 건강나이 산출 및 임상의학적 근거
    # ==================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, COLOR_LIGHT_BG)
    add_slide_header(s5, "통합 건강나이(3BODY Core Balance Age) 산출 및 임상 근거", "Integrated Biological Age Model")

    # 좌측: 통합 나이 산출 아키텍처
    create_card_shape(s5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_l5 = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_l5 = tb_l5.text_frame
    tf_l5.word_wrap = True
    tf_l5.paragraphs[0].text = "🎛️ 통합 건강나이 산출 프레임워크"
    tf_l5.paragraphs[0].font.name = FONT_NAME
    tf_l5.paragraphs[0].font.size = Pt(16)
    tf_l5.paragraphs[0].font.bold = True
    tf_l5.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_l5.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_l5, "• 다차원 가중치 합산 역산 알고리즘", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l5, "몸나이(40%), 뇌나이(20%), 마음나이(25%), 얼굴나이(15%)를 각 시스템 기여도에 따라 가중 결합하여 통합 '3BODY 밸런스 건강나이'를 산출합니다. 이는 단일 수치로 건강 상태를 종합 직관할 수 있게 합니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_l5, "• 100점 만점 절대 채점과 노화 속도 정밀 보간", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l5, "전체 검사 데이터를 100점 기준 절대 점수화하고 70점을 실제 나이(0 증감)의 기준점으로 삼아, 1점 증감당 0.4세 편차로 생물학적 노화 가속/감속 속도를 환산해 냅니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # 우측: PNI 연쇄 작용과 임상적 의의
    create_card_shape(s5, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_r5 = s5.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_r5 = tb_r5.text_frame
    tf_r5.word_wrap = True
    tf_r5.paragraphs[0].text = "🔬 PNI 및 신체-인지 노화 양방향 연쇄 기전"
    tf_r5.paragraphs[0].font.name = FONT_NAME
    tf_r5.paragraphs[0].font.size = Pt(16)
    tf_r5.paragraphs[0].font.bold = True
    tf_r5.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_r5.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_r5, "• 몸-마음-뇌의 생리학적 유기성 (HPA 축)", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r5, "심리적 스트레스(마음)는 시상하부-뇌하수체-부신 축을 통해 코르티솔을 과분비시켜 신체 염증(몸)을 초래하고, 이는 염증성 사이토카인을 통해 전두엽(뇌)의 인지 기능 및 정서 제어력을 떨어뜨려 연쇄 노화를 유도합니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_r5, "• 통합 진단의 임상적 우수성", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r5, "특정 영역만 관리해서는 도미노처럼 연결된 만성 노화를 막을 수 없습니다. 3BODY 통합 나이는 몸, 마음, 뇌를 다차원적으로 아우르는 후성유전학적 건강 진단의 표준을 제시합니다.", 11, False, COLOR_TEXT_MUTED, 6)


    # ==================================================
    # SLIDE 6: 7CODE 에너지 시스템의 신경/내분비학 매핑
    # ==================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, COLOR_LIGHT_BG)
    add_slide_header(s6, "7CODE 에너지 시스템의 자율신경총 및 내분비선 매핑", "7Code Integrated Model")

    # 설명 텍스트
    desc_box = s6.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.6))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.paragraphs[0].text = "7CODE 시스템은 동양의 차크라(Chakra) 에너지 개념을 현대 해부생리학의 자율신경총 및 내분비선 작용과 비교 매핑하여 과학적 설득력을 극대화한 설명 모델입니다."
    tf_desc.paragraphs[0].font.name = FONT_NAME
    tf_desc.paragraphs[0].font.size = Pt(11)
    tf_desc.paragraphs[0].font.bold = False
    tf_desc.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED

    # 테이블 생성 (8행 5열)
    t_shape = s6.shapes.add_table(8, 5, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.7))
    table = t_shape.table
    table.columns[0].width = Inches(1.8)  # 코드 구분
    table.columns[1].width = Inches(2.3)  # 자율신경총
    table.columns[2].width = Inches(2.2)  # 주요 내분비선
    table.columns[3].width = Inches(4.8)  # 조절 기능 (생리적/심리적)
    table.columns[4].width = Inches(1.0)  # 학술설명

    headers = ["코드 구분", "대응 자율신경총", "대응 내분비선", "주요 기능 (생리적/심리적)", "학술근거"]
    for i, h in enumerate(headers):
        format_cell_content(table.cell(0, i), h, 11, True, COLOR_WHITE)
    style_table_header(table, 5)

    data_7code = [
        ("1CODE (기초/생존)", "골반 신경총 (Pelvic Plexus)", "부신 (Adrenal Gland)", "코르티솔 등 스트레스 호르몬 통제, 그라운딩, 기초 체력", "생리학적"),
        ("2CODE (감정/관계)", "하복부/요천추 신경총", "생식선 (Ovaries/Testes)", "에너지와 체액의 순환 작용, 감정 유연성, 창조력", "상관성"),
        ("3CODE (의욕/추진)", "복강/태양신경총 (Celiac Plexus)", "췌장 (Pancreas)", "대사 작용, 자율신경계 조절, 자기 통제 및 추진력", "해부학적"),
        ("4CODE (사랑/안정)", "심장 신경총 (Cardiac Plexus)", "흉선 (Thymus Gland)", "면역 T세포 분비 도움, 심박수 조절, 정서적 안정 및 치유", "대조"),
        ("5CODE (소통/표현)", "경부/인두 신경총 (Cervical Plexus)", "갑상선 (Thyroid Gland)", "대사 속도 제어, 자기 소통 및 언어 표현, 목/어깨 이완", "매핑"),
        ("6CODE (집중/통찰)", "두개강 내 신경망 (송과체 주변)", "뇌하수체 (Pituitary Gland)", "멜라토닌 수면 주기 통제, 직관력, 시각 처리와 집중력", "과학적"),
        ("7CODE (방향/의미)", "대뇌피질 및 전반 중추신경계", "송과체 (Pineal Gland)", "고위 인지 기능 조절, 통합적 평온함, 인생 목표 정립", "타당성확보")
    ]

    for idx, row_data in enumerate(data_7code, 1):
        bg_col = COLOR_WHITE if idx % 2 != 0 else COLOR_LIGHT_BLUE
        for col_idx, text in enumerate(row_data):
            cell = table.cell(idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_col
            
            align = PP_ALIGN.CENTER if col_idx in [0, 1, 2, 4] else PP_ALIGN.LEFT
            is_bold = col_idx == 0
            text_color = COLOR_PRIMARY if col_idx == 0 else COLOR_TEXT_DARK
            format_cell_content(cell, text, 10, is_bold, text_color, align)


    # ==================================================
    # SLIDE 7: 학술적 근거 및 과학적 기전
    # ==================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, COLOR_LIGHT_BG)
    add_slide_header(s7, "통합 웰니스의 신경과학적 규명 및 임상 연구 성과", "Holistic Wellness Clinical Studies")

    # 좌측: 워싱턴대 SCAN 연구 (Nature, 2023)
    create_card_shape(s7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_l7 = s7.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_l7 = tb_l7.text_frame
    tf_l7.word_wrap = True
    tf_l7.paragraphs[0].text = "🔬 워싱턴 대학교 SCAN 연구 (Nature, 2023)"
    tf_l7.paragraphs[0].font.name = FONT_NAME
    tf_l7.paragraphs[0].font.size = Pt(16)
    tf_l7.paragraphs[0].font.bold = True
    tf_l7.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_l7.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_l7, "• 신체-인지 작용 네트워크(SCAN)의 발견", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l7, "대뇌 1차 운동 피질(Motor Cortex) 내부에 신체 움직임, 목표 지향적 계획, 자율신경계 조절(혈압, 심박 등) 기능이 완전히 얽혀 하나의 연동 네트워크로 작용함을 밝혀냈습니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_l7, "• 3BODY 통합 치료의 해부학적 단초", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l7, "특정 호흡이나 움직임(Body)이 어떻게 인지 제어(Brain)와 스트레스 조절(Mind)에 직접적인 피드백 경로를 만들어내는지 물리적으로 증명한 우수한 모델입니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # 우측: 뇌교육 명상(BEM)의 임상적 실효성 및 유전자 바이오마커
    create_card_shape(s7, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_r7 = s7.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_r7 = tb_r7.text_frame
    tf_r7.word_wrap = True
    tf_r7.paragraphs[0].text = "🩸 뇌교육 명상(BEM) 임상 및 염증 유전자 변화"
    tf_r7.paragraphs[0].font.name = FONT_NAME
    tf_r7.paragraphs[0].font.size = Pt(16)
    tf_r7.paragraphs[0].font.bold = True
    tf_r7.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_r7.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_r7, "• 만성 염증 유전자 발현량 억제 (후성유전학적 영향)", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r7, "BEM 훈련 8주 만에 스트레스와 세포 노화의 원인인 염증 유전자 발현이 유의하게 감소했습니다 (NFKB2 70% 감소, RELA 50% 감소, IL1B 80% 급감).", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_r7, "• 심혈관 위험 인자 개선 및 대사 개선", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r7, "무작위 대조군 연구 결과, BEM 훈련을 수행한 그룹에서 저밀도 지단백(LDL) 콜레스테롤 수치가 평균 13.82 mg/dL 감소하는 등 신진대사의 객관적 개선을 입증했습니다.", 11, False, COLOR_TEXT_MUTED, 6)


    # ==================================================
    # SLIDE 8: 습관 형성 3단계 프로그램의 과학적 기전
    # ==================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, COLOR_LIGHT_BG)
    add_slide_header(s8, "21일 · 66일 · 100일 습관 형성의 뇌과학적 원리", "Habit Formation & Neuroplasticity")

    card_w = Inches(3.6)
    card_h = Inches(5.1)

    # 1단계: 21일 적응기
    create_card_shape(s8, Inches(0.8), Inches(1.5), card_w, card_h)
    tb1 = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "🌱 21일 적응기 (Adaptation)"
    tf1.paragraphs[0].font.name = FONT_NAME
    tf1.paragraphs[0].font.size = Pt(15)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    tf1.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf1, "• 뇌의 새로운 자극 적응 단계", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf1, "맥스웰 몰츠 박사의 저술에서 유래한 ‘21일 법칙’은 뇌가 낯선 행동이나 자아 이미지를 저항 없이 수용하기 시작하는 최소한의 적응 기간을 의미합니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf1, "• 뇌세포 화학 신호의 초기 점화", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf1, "명상과 에너지 조절 훈련에 대한 초기 신경 경로의 생화학적 반응을 활성화시키는 진입 프로그램 설계 기준입니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2단계: 66일 자동화기
    create_card_shape(s8, Inches(4.8), Inches(1.5), card_w, card_h, COLOR_WHITE, COLOR_SECONDARY)
    tb2 = s8.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "🌿 66일 자동화기 (Automaticity)"
    tf2.paragraphs[0].font.name = FONT_NAME
    tf2.paragraphs[0].font.size = Pt(15)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf2.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf2, "• UCL 필리파 랠리 교수팀 연구", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf2, "새로운 행동이 의식적 노력 없이 자동화(Automaticity)되는 데 평균 66일이 걸린다고 보고했습니다 (개인 및 난이도별 18~254일 소요).", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf2, "• 안정적 항상성 형성 구간", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf2, "훈련을 무의식적 일상 루틴으로 정착시켜, 좋아졌다 무너지는 일시적 적응의 한계를 넘는 전략적 기간으로 활용합니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3단계: 100일 수초화기
    create_card_shape(s8, Inches(8.8), Inches(1.5), Inches(3.7), card_h)
    tb3 = s8.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "🌳 100일 뇌 재배선 (Myelination)"
    tf3.paragraphs[0].font.name = FONT_NAME
    tf3.paragraphs[0].font.size = Pt(15)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_ACCENT
    tf3.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf3, "• 신경가소성 및 수초화 원리 참고", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf3, "반복적인 자극은 신경의 정보 처리 속도를 가속하는 축삭 수초화(Myelination)를 유도하여 신경 연결을 강화하는 데 기여합니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf3, "• 100일의 실천적 목표 설계", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf3, "뉴런 결속을 유도하는 헤브의 법칙(Hebbian learning)을 반영하여, 삶의 근본적인 인지 필터와 태도를 리모델링하는 실천적 기간입니다.", 10.5, False, COLOR_TEXT_MUTED, 6)


    # ==================================================
    # SLIDE 9: 대규모 성과 및 공공/기업 도입 레퍼런스
    # ==================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, COLOR_LIGHT_BG)
    add_slide_header(s9, "대규모 점검 성과 및 신뢰도 검증 레퍼런스", "References & Track Record")

    # 1. 성과 요약 카드 (좌측)
    create_card_shape(s9, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.1))
    tb_stats = s9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(4.1), Inches(4.7))
    tf_stats = tb_stats.text_frame
    tf_stats.word_wrap = True
    tf_stats.paragraphs[0].text = "🚀 단기 대규모 측정 성과"
    tf_stats.paragraphs[0].font.name = FONT_NAME
    tf_stats.paragraphs[0].font.size = Pt(16)
    tf_stats.paragraphs[0].font.bold = True
    tf_stats.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_stats.paragraphs[0].space_after = Pt(20)

    add_paragraph_to_tf(tf_stats, "출시 2달 미만 누적 측정자", 12, True, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_stats, "8,000명 돌파", 32, True, COLOR_ACCENT, 16)
    
    add_paragraph_to_tf(tf_stats, "“단 2개월 만에 8,000명 이상의 고유 점검 데이터를 안정적으로 누적했습니다. 이는 하드웨어 리소스 최적화 및 2D 컴퓨터 비전 실시간 연산 안전성이 실전 현장에서 완벽하게 검증되었음을 대변합니다.”", 11, False, COLOR_TEXT_DARK, 6)

    # 2. 공공기관 및 기업 레퍼런스 카드 (우측)
    create_card_shape(s9, Inches(5.6), Inches(1.5), Inches(6.9), Inches(5.1))
    tb_ref = s9.shapes.add_textbox(Inches(5.8), Inches(1.7), Inches(6.5), Inches(4.7))
    tf_ref = tb_ref.text_frame
    tf_ref.word_wrap = True
    tf_ref.paragraphs[0].text = "🏢 관공서 및 대기업 · 단체 도입 레퍼런스"
    tf_ref.paragraphs[0].font.name = FONT_NAME
    tf_ref.paragraphs[0].font.size = Pt(16)
    tf_ref.paragraphs[0].font.bold = True
    tf_ref.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_ref.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_ref, "• 주요 지자체 복지 사업 지원", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "춘천시청, 포항시청, 춘천교육청, 천안시 행복키움단 등 시청 및 교육 기관 산하 건강/복지 프로그램으로 도입 및 운영되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_ref, "• 전문 의료진 및 체육 · 공익 단체", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "서울약사협회, 대구약사협회, 경북약사협회, 평택체육회, 안동노인회 등 높은 객관성이 요구되는 약사 및 보건 체육 단체 회원들을 대상으로 성공적인 단체 건강 지표 수립 검사를 실행했습니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_ref, "• 대기업 복지 연계", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "IT 대기업인 카카오 등 소속 임직원의 멘탈 웰니스 측정 및 직장인 증후군 자가 예방 지원 툴로써 AI 코드맵 점검이 기여하고 있습니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # PPT 파일 저장
    output_ppt = '3바디_7코드_AI코드맵_신뢰성_홍보브로셔.pptx'
    try:
        prs.save(output_ppt)
        print(f"[완료] AI코드맵 신뢰성 홍보브로셔 PPTX 생성 성공: {output_ppt}")
    except PermissionError:
        print(f"[오류] {output_ppt} 파일이 현재 열려 있거나 사용 중입니다.")

if __name__ == '__main__':
    generate_presentation()
