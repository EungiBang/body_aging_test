# 부경대학교 3바디7코드 AI 코드맵 팩트체크 및 일반인 공감 슬라이드를 추가 반영하여 최종 PPTX를 생성하는 스크립트
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 테마 색상 (신뢰감을 주는 프리미엄 딥 네이비 & 아일랜드 블루 테마)
COLOR_DARK_BG = RGBColor(15, 23, 42)      # 슬레이트 블랙/딥 네이비
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # 은은한 그레이시 화이트
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PRIMARY = RGBColor(27, 54, 93)      # 다크 네이비
COLOR_SECONDARY = RGBColor(59, 130, 246)  # 아일랜드 블루
COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # 차콜 블랙
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 뮤트 그레이
COLOR_ACCENT = RGBColor(245, 158, 11)     # 골드/오렌지 액센트
COLOR_CARD_BORDER = RGBColor(226, 232, 240)
COLOR_LIGHT_BLUE = RGBColor(239, 246, 255)

FONT_NAME = "Malgun Gothic"

def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title, subtitle):
    # 상단 대분류/서브 타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.3))
    p_sub = sub_box.text_frame.paragraphs[0]
    p_sub.text = subtitle.upper()
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(9.5)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.6))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_DARK

def create_card_shape(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER):
    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = bg_color
    card_shape.line.color.rgb = border_color
    card_shape.line.width = Pt(1)
    return card_shape

def add_paragraph_to_tf(tf, text, font_size=11, is_bold=False, font_color=COLOR_TEXT_DARK, space_after=6, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    return p

def format_cell_content(cell, text, font_size=10, is_bold=False, font_color=COLOR_TEXT_DARK, alignment=PP_ALIGN.CENTER):
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color

def style_table_header(table, col_count):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

def generate_pukyong_presentation():
    # PPT 생성 시작 (Widescreen 16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==================================================
    # SLIDE 1: 표지 슬라이드
    # ==================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_DARK_BG)

    left_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_SECONDARY
    left_bar.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(2.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_SECONDARY
    badge.line.fill.background()
    badge.text_frame.text = "통합 건강 패러다임 제안"
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = FONT_NAME
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_WHITE

    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "부경대학교 교직원 및 학생 대상 복지·건강 증진 프로그램"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(16)
    p0.font.color.rgb = COLOR_SECONDARY
    p0.font.bold = True
    p0.space_after = Pt(12)

    p1 = tf.add_paragraph()
    p1.text = "3바디 7코드 AI 코드맵 학술적 근거 보고서"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "몸·마음·뇌의 통합적 균형과 7코드 에너지 기반 웰니스 진단 시스템의 타당성 분석"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_WHITE
    p2.font.bold = False
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "공동 연구 개발: 국가공인 브레인트레이너 전문가 단체 · 한국뇌과학연구원 · 글로벌사이버대학교"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(6)

    p4 = tf.add_paragraph()
    p4.text = "제공 기관: 브레인트레이닝센터"
    p4.font.name = FONT_NAME
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_TEXT_MUTED

    # ==================================================
    # SLIDE 2: [NEW] 왜 몸·마음·뇌의 통합적 건강인가?
    # ==================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_LIGHT_BG)
    add_slide_header(s2, "건강의 새로운 관점: 신체·마음·뇌의 통합적 균형", "Integrated Health Needs")

    card_w2 = Inches(3.6)
    card_h2 = Inches(5.1)

    # 1. 몸 (신체 건강에 치중된 관리의 한계)
    create_card_shape(s2, Inches(0.8), Inches(1.5), card_w2, card_h2)
    tb_b1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b1 = tb_b1.text_frame
    tf_b1.word_wrap = True
    tf_b1.paragraphs[0].text = "💪 신체 위주 관리의 한계"
    tf_b1.paragraphs[0].font.name = FONT_NAME
    tf_b1.paragraphs[0].font.size = Pt(15)
    tf_b1.paragraphs[0].font.bold = True
    tf_b1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b1.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b1, "• 반쪽짜리 건강 관리", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b1, "대부분의 사람들은 건강이라고 하면 주로 헬스, 요가, 식단 등 육체적인 관리에만 전념합니다. 하지만 몸만 가꾼다고 해서 삶의 활력이 저절로 채워지지는 않습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b1, "• 유기적 연관성", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b1, "몸의 긴장은 뇌의 스트레스를 유발하고, 지친 뇌는 다시 신체 소화기계 및 자율신경계 저하로 이어집니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2. 마음 (소홀하기 쉬운 마음의 병)
    create_card_shape(s2, Inches(4.8), Inches(1.5), card_w2, card_h2)
    tb_b2 = s2.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b2 = tb_b2.text_frame
    tf_b2.word_wrap = True
    tf_b2.paragraphs[0].text = "❤️ 방치되는 마음의 피로와 스트레스"
    tf_b2.paragraphs[0].font.name = FONT_NAME
    tf_b2.paragraphs[0].font.size = Pt(15)
    tf_b2.paragraphs[0].font.bold = True
    tf_b2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b2.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b2, "• 일상 스트레스 인지율 25.9% (질병관리청)", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b2, "성인 4명 중 1명은 일상 중 심한 스트레스를 느끼며, 특히 주 경제 활동층인 40대(35.1%)와 30대(34.7%)의 스트레스 인지가 매우 높게 나타납니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b2, "• 직장인 번아웃 경험률 60~70% 돌파", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b2, "직무 및 대인관계 스트레스로 인해 정서적 탈진 상태를 경험한 비율은 직장인 대상 여러 통계 조사에서 최소 60%에서 최대 90% 이상으로 집계되었습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3. 뇌 (우려와 대책의 공백)
    create_card_shape(s2, Inches(8.8), Inches(1.5), Inches(3.7), card_h2)
    tb_b3 = s2.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf_b3 = tb_b3.text_frame
    tf_b3.word_wrap = True
    tf_b3.paragraphs[0].text = "🧠 치매 걱정과 관리 방법의 부재"
    tf_b3.paragraphs[0].font.name = FONT_NAME
    tf_b3.paragraphs[0].font.size = Pt(15)
    tf_b3.paragraphs[0].font.bold = True
    tf_b3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b3.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b3, "• 가장 두려운 노후 질병 1위 치매 (43%)", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b3, "심평원 설문에 따르면 노인들이 가장 걱정하는 질환은 암(33%)보다도 치매(43%)가 압도적 1위로 나타나 강한 우려감을 대변합니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b3, "• 경도인지장애 유병률 28.4% (중앙치매센터)", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b3, "65세 이상 노인의 치매 유병률은 9.25%로 고령층 10명 중 1명에 달하며, 전 단계인 경도인지장애 유병률은 28.4%에 육박하나 일상의 체계적 두뇌 훈련이 부재한 상태입니다.", 10.5, False, COLOR_TEXT_MUTED, 8)

    # ==================================================
    # SLIDE 3: [NEW] 컨디션을 결정하는 열쇠: 7코드 에너지 관리
    # ==================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_LIGHT_BG)
    add_slide_header(s3, "컨디션과 건강을 지배하는 에너지 관리의 중요성", "Energy & 7Code Concept")

    # 1. 에너지 상태에 따른 선순환 vs 악순환 비교 카드 (좌측)
    create_card_shape(s3, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_el = s3.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_el = tb_el.text_frame
    tf_el.word_wrap = True
    tf_el.paragraphs[0].text = "🔄 컨디션을 결정하는 생체 에너지 흐름"
    tf_el.paragraphs[0].font.name = FONT_NAME
    tf_el.paragraphs[0].font.size = Pt(16)
    tf_el.paragraphs[0].font.bold = True
    tf_el.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_el.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_el, "• 에너지가 떨어질 때 (악순환)", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_el, "체내 에너지가 방전되면 의욕 저하가 오고, 감정 기복이 심해집니다. 이 감정은 자율신경을 거쳐 신체 피로를 만들며, 결국 뇌의 스트레스 과부하를 부릅니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    
    add_paragraph_to_tf(tf_el, "• 에너지가 충전될 때 (선순환)", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_el, "에너지가 흐르고 중심이 잡히면 즉각 마음의 안정과 기분 개선이 옵니다. 신체 대사 기능과 자가 치유력이 증가하며 뇌파가 안정되어 판단력과 집중력이 향상됩니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2. 고대의 지혜와 현대 과학의 연결 (우측)
    create_card_shape(s3, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_er = s3.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_er = tb_er.text_frame
    tf_er.word_wrap = True
    tf_er.paragraphs[0].text = "📜 고대의 건강 지혜에서 현대 7코드로의 진화"
    tf_er.paragraphs[0].font.name = FONT_NAME
    tf_er.paragraphs[0].font.size = Pt(16)
    tf_er.paragraphs[0].font.bold = True
    tf_er.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_er.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_er, "• 동양 철학의 '차크라'와 한국의 '삼단전'", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_er, "고대의 현인들은 건강을 에너지 관점으로 바라보았습니다. 인도 요가 철학의 7개 차크라, 그리고 한국 전통 건강법인 3개 단전(하단전, 중단전, 상단전) 센터 이론이 그 대표적인 지혜입니다.", 10.5, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_er, "• 현대 과학으로 재해석한 7CODE", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_er, "브레인트레이닝센터는 이 고대 신체 에너지 센터들의 실체를 현대 해부생리학의 자율신경총 및 호르몬 내분비선 위치와 연계했습니다. 일반인 누구나 자신의 무너진 컨디션 상태를 쉽게 파악하고 공감할 수 있도록 '7CODE'로 명료하게 계량화하여 풀었습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 4: 공동 연구 개발 및 학술적 토대
    # ==================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, COLOR_LIGHT_BG)
    add_slide_header(s4, "3바디 7코드 AI 코드맵의 공동 연구 개발진", "R&D Background")

    # 3개 기관 소개 카드 배치
    card_w = Inches(3.6)
    card_h = Inches(5.1)
    
    # 1. 한국뇌과학연구원 (KIBS)
    create_card_shape(s4, Inches(0.8), Inches(1.5), card_w, card_h)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🏢 한국뇌과학연구원 (KIBS)"
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf, "• 뇌교육 원천기술 연구 개발", 11, True, COLOR_TEXT_DARK, 8)
    add_paragraph_to_tf(tf, "인간 뇌의 잠재력을 주체적으로 활용하기 위한 기초 연구와 학술적 메커니즘을 제안하는 뇌교육 특화 연구 기관입니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf, "• 뇌운영시스템 (B.O.S) 제공", 11, True, COLOR_TEXT_DARK, 8)
    add_paragraph_to_tf(tf, "뇌를 스스로 주체적으로 관리하고 조절하는 'B.O.S 5단계' 프레임워크를 제안하여 프로그램의 뇌과학적 훈련 방법론을 뒷받침합니다.", 10.5, False, COLOR_TEXT_MUTED, 8)

    # 2. 글로벌사이버대학교
    create_card_shape(s4, Inches(4.8), Inches(1.5), card_w, card_h)
    tb = s4.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🎓 글로벌사이버대학교"
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf, "• 명상치유학 및 뇌교육 연구진", 11, True, COLOR_TEXT_DARK, 8)
    add_paragraph_to_tf(tf, "명상치유 및 뇌교육 융합 전공 학과 교수진이 프로그램의 교육학적·심리학적 타당성을 검토하고 학문적 융합 웰니스 이론을 접목했습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf, "• 뇌기반 교육의 실천적 성과", 11, True, COLOR_TEXT_DARK, 8)
    add_paragraph_to_tf(tf, "현대 교육 및 복지 현장에서 주의집중력, 신체 조절력, 정서 관리를 통합 연계하는 실천적 교육 모델을 제공합니다.", 10.5, False, COLOR_TEXT_MUTED, 8)

    # 3. 국가공인 브레인트레이너 전문가 단체
    create_card_shape(s4, Inches(8.8), Inches(1.5), Inches(3.7), card_h)
    tb = s4.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🏅 브레인트레이너 협회"
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf, "• 국가공인 두뇌훈련 자격 단체", 11, True, COLOR_TEXT_DARK, 8)
    add_paragraph_to_tf(tf, "교육부 공인 두뇌 관리 및 훈련 자격을 보유한 현장 트레이너들이 직접 프로그램 실무 설계 및 검사 피드백 구축에 참여했습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf, "• 축적된 현장 적용 데이터", 11, True, COLOR_TEXT_DARK, 8)
    add_paragraph_to_tf(tf, "수많은 현장 상담 데이터에 기반해 직관적이고 자발적으로 자세 및 스트레스 상태를 교정할 수 있도록 맞춤형 가이드를 설계했습니다.", 10.5, False, COLOR_TEXT_MUTED, 8)

    # ==================================================
    # SLIDE 5: AI 측정 기술 및 특허출원 신청 현황
    # ==================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, COLOR_LIGHT_BG)
    add_slide_header(s5, "다차원 AI 정밀 측정 핵심 기술 및 특허출원 신청", "AI Technology & Patent")

    # 상단 요약 설명 박스
    create_card_shape(s5, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1), COLOR_LIGHT_BLUE, COLOR_SECONDARY)
    tb_top = s5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.9))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = True
    p_top = tf_top.paragraphs[0]
    p_top.text = "💡 독자적 AI 컴퓨터 비전 기반 건강 측정 엔진 특허출원 신청 진행 중"
    p_top.font.name = FONT_NAME
    p_top.font.size = Pt(13)
    p_top.font.bold = True
    p_top.font.color.rgb = COLOR_PRIMARY
    p_top.space_after = Pt(4)
    add_paragraph_to_tf(tf_top, "현장 간이 측정의 부정확성을 보장하고 기술 신뢰도를 강화하기 위해, AI 컴퓨터 비전 및 실시간 물리 트래킹 기술에 대해 관련 특허 출원 절차를 밟고 있습니다. 기관 도입 신뢰성을 명확하게 제시합니다.", 10.5, False, COLOR_TEXT_DARK, 0)

    # 하단 3단 상세 기술 카드
    card_y = Inches(2.7)
    card_h2 = Inches(4.0)
    card_w2 = Inches(3.6)

    # 기술 1: 실시간 부정방지(Cheat Prevention) 기술
    create_card_shape(s5, Inches(0.8), card_y, card_w2, card_h2)
    tb1 = s5.shapes.add_textbox(Inches(1.0), card_y + Inches(0.2), Inches(3.2), card_h2 - Inches(0.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "🛡️ 실시간 자세 흔들림 감지"
    tf1.paragraphs[0].font.name = FONT_NAME
    tf1.paragraphs[0].font.size = Pt(14)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf1.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf1, "• 미세 흔들림 추적 (Sway Score)", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf1, "코와 어깨 좌표의 실시간 움직임을 정규화하여 밸런스 점수를 수치화하며, 일시적 정지 자세의 왜곡을 방지합니다.", 10, False, COLOR_TEXT_MUTED, 8)
    add_paragraph_to_tf(tf1, "• 실시간 발 디딤 (Foot Drop) 센싱", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf1, "양 발목 Y축 좌표의 프레임별 변동을 연산하여, 한발 서기 중 발이 바닥에 닿는 행위를 자동 판별합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 기술 2: WebGL 가속 및 오클루전 극복
    create_card_shape(s5, Inches(4.8), card_y, card_w2, card_h2)
    tb2 = s5.shapes.add_textbox(Inches(5.0), card_y + Inches(0.2), Inches(3.2), card_h2 - Inches(0.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "⚡ 가속화 및 오클루전 극복"
    tf2.paragraphs[0].font.name = FONT_NAME
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf2.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf2, "• VRAM 메모리 관리 최적화", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf2, "웹 환경에서 매 프레임별 WebGL 텍스처를 파기하도록 설계하여, 컴퓨터 메모리 과부하 및 다운 현상을 근본적으로 해결했습니다.", 10, False, COLOR_TEXT_MUTED, 8)
    add_paragraph_to_tf(tf2, "• 관절 겹침(Occlusion) 자동 보정", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf2, "측면 스쿼트와 대각선 푸시업 구도에서 가시성이 확보된 최적의 관절 좌표를 선별하여 트래킹 정확도를 높입니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 기술 3: 가중치 기반 신체나이 환산 알고리즘
    create_card_shape(s5, Inches(8.8), card_y, Inches(3.7), card_h2)
    tb3 = s5.shapes.add_textbox(Inches(9.0), card_y + Inches(0.2), Inches(3.3), card_h2 - Inches(0.4))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "📊 신체/체력 나이 역산 알고리즘"
    tf3.paragraphs[0].font.name = FONT_NAME
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf3.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf3, "• 데이터 기반 절대평가 체점", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf3, "성별 및 연령대별 평균 달성 표준에 따라 선형 보간 채점하여, AI의 주관적이고 후한 점수 평가 편향(인플레)을 차단합니다.", 10, False, COLOR_TEXT_MUTED, 8)
    add_paragraph_to_tf(tf3, "• 다차원 가중치 합산 공식", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf3, "근력, 균형, 자세, 유연성 등 각 테스트 항목별로 적정 가중치 비율을 설정한 후 최종 건강 연령 지표를 도출합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 6: 3BODY(몸·마음·뇌) 통합 구조의 학술적 근거
    # ==================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, COLOR_LIGHT_BG)
    add_slide_header(s6, "3BODY(몸·마음·뇌) 연계의 현대 의과학적 접점", "3Body Scientific Reference")

    # 좌측: 심신신경면역학(PNI) 근거 카드
    create_card_shape(s6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_l = s6.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "🧠 심신신경면역학(PNI) 및 HPA 축"
    tf_l.paragraphs[0].font.name = FONT_NAME
    tf_l.paragraphs[0].font.size = Pt(16)
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_l.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_l, "• 마음과 몸의 쌍방향적 상호작용 지지", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l, "스트레스, 불안 등의 심리적 상태는 시상하부-뇌하수체-부신 축(HPA Axis)을 자극하고 코르티솔 분비를 유도하여 신체에 직접적인 피드백을 전달할 수 있습니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_l, "• 신체 피로와 인지 저하의 양방향 연관성", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l, "신체 내부의 염증 물질(사이토카인 등)이 미주신경 등을 거쳐 뇌에 전달됨으로써 기분 저하, 무기력감, 주의집중력의 변화(Sickness Behavior)를 일으키는 접점을 설명합니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # 우측: 워싱턴대 SCAN 연구 (Nature, 2023)
    create_card_shape(s6, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_r = s6.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "🔬 워싱턴 대학교 SCAN 연구 (Nature, 2023)"
    tf_r.paragraphs[0].font.name = FONT_NAME
    tf_r.paragraphs[0].font.size = Pt(16)
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_r.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_r, "• SCAN(Somato-Cognitive Action Network) 보고", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r, "Nature 및 후속 연구에 따르면, 대뇌 1차 운동 피질 내부에 신체 움직임, 목표 지향 행동, 생리 상태(혈압 등)를 통합하는 기능적 네트워크의 존재가 보고되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_r, "• 3BODY 통합의 신경학적 참고 근거", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r, "이러한 SCAN의 구조는 몸의 움직임과 인지·생리 기능이 긴밀히 연관되어 있음을 시사하며, 몸·마음·뇌를 통합적으로 다루는 3BODY 접근의 접점을 뒷받침하는 우수한 과학적 참고 모델입니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 7: 7CODE 에너지 시스템의 신경/내분비학 매핑
    # ==================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, COLOR_LIGHT_BG)
    add_slide_header(s7, "7CODE 에너지 시스템의 자율신경총 및 내분비선 해석", "7Code Integrated Model")

    # 설명 텍스트
    desc_box = s7.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.6))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.paragraphs[0].text = "7CODE 시스템은 전통적인 7개 에너지 센터(차크라) 개념을 현대 해부생리학의 주요 자율신경총 및 내분비선과 비교·해석한 설명 모델입니다."
    tf_desc.paragraphs[0].font.name = FONT_NAME
    tf_desc.paragraphs[0].font.size = Pt(11)
    tf_desc.paragraphs[0].font.bold = False
    tf_desc.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED

    # 테이블 생성 (8행 5열)
    t_shape = s7.shapes.add_table(8, 5, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.7))
    table = t_shape.table
    table.columns[0].width = Inches(1.8)  # 코드 구분
    table.columns[1].width = Inches(2.3)  # 자율신경총
    table.columns[2].width = Inches(2.2)  # 주요 내분비선
    table.columns[3].width = Inches(4.8)  # 조절 기능 (생리적/심리적)
    table.columns[4].width = Inches(1.0)  # 신뢰성 근거

    headers = ["코드 구분", "대응 자율신경총", "대응 내분비선", "주요 기능 (생리적/심리적)", "학술설명"]
    for i, h in enumerate(headers):
        format_cell_content(table.cell(0, i), h, 11, True, COLOR_WHITE)
    style_table_header(table, 5)

    data_7code = [
        ("1CODE (기초/생존)", "골반 신경총 (Pelvic Plexus)", "부신 (Adrenal Gland)", "스트레스 호르몬 조절, 투쟁-도피 반응, 기초 체력 및 안정감", "생리학적"),
        ("2CODE (감정/관계)", "하복부/요천추 신경총", "생식선 (Ovaries/Testes)", "체액 및 에너지 순환, 감정 흐름, 수용력 및 창조력", "상관성을"),
        ("3CODE (의욕/추진)", "복강/태양신경총", "췌장 (Pancreas)", "소화 및 대사 작용, 자율신경 조절, 자존감 및 통제력", "이해하기"),
        ("4CODE (사랑/안정)", "심장 신경총 (Cardiac Plexus)", "흉선 (Thymus Gland)", "면역계 조절 도움, 심박수 조절, 정서적 안정 및 회복력", "위한"),
        ("5CODE (소통/표현)", "경부/인두 신경총", "갑상선 (Thyroid Gland)", "신진대사 조절, 언어 표현 및 소통, 목/어깨 근육 이완", "통합적"),
        ("6CODE (집중/통찰)", "두개강 내 신경망 (송과체 주변)", "뇌하수체 (Pituitary Gland)", "멜라토닌(수면 주기) 및 전체 호르몬 통제, 직관, 집중력", "참고적"),
        ("7CODE (방향/의미)", "대뇌피질 및 중추신경계", "송과체 (Pineal Gland)", "고위 인지 기능 조절, 삶의 의미 및 정체성 방향성 수립", "설명 모델")
    ]

    for idx, row_data in enumerate(data_7code, 1):
        bg_col = COLOR_WHITE if idx % 2 != 0 else COLOR_LIGHT_BLUE
        for col_idx, text in enumerate(row_data):
            cell = table.cell(idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_col
            
            align = PP_ALIGN.CENTER if col_idx in [0, 1, 2, 4] else PP_ALIGN.LEFT
            is_bold = col_idx == 0
            text_color = COLOR_PRIMARY if col_idx == 0 else COLOR_TEXT_DARK
            format_cell_content(cell, text, 10, is_bold, text_color, align)

    # ==================================================
    # SLIDE 8: 습관 형성 3단계 프로그램의 과학적 기전
    # ==================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, COLOR_LIGHT_BG)
    add_slide_header(s8, "21일 · 66일 · 100일 습관 형성의 뇌과학적 원리", "Habit Formation & Neuroplasticity")

    card_w = Inches(3.6)
    card_h = Inches(5.1)

    # 1단계: 21일 적응기
    create_card_shape(s8, Inches(0.8), Inches(1.5), card_w, card_h)
    tb1 = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "🌱 21일 적응기 (Adaptation)"
    tf1.paragraphs[0].font.name = FONT_NAME
    tf1.paragraphs[0].font.size = Pt(15)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    tf1.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf1, "• 뇌의 새로운 자극 적응 단계", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf1, "맥스웰 몰츠 박사의 저술에서 유래한 ‘21일 법칙’은 뇌가 낯선 행동이나 자아 이미지를 저항 없이 수용하기 시작하는 최소한의 적응 기간을 의미합니다. (경험적 통찰의 참고치)", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf1, "• 뇌세포 화학 신호의 초기 점화", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf1, "명상과 에너지 조절 훈련에 대한 초기 신경 경로의 생화학적 반응을 활성화시키는 진입 프로그램 설계 기준입니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2단계: 66일 자동화기
    create_card_shape(s8, Inches(4.8), Inches(1.5), card_w, card_h, COLOR_WHITE, COLOR_SECONDARY)
    tb2 = s8.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "🌿 66일 자동화기 (Automaticity)"
    tf2.paragraphs[0].font.name = FONT_NAME
    tf2.paragraphs[0].font.size = Pt(15)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf2.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf2, "• UCL 필리파 랠리 교수팀 연구", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf2, "새로운 행동이 의식적 노력 없이 자동화(Automaticity)되는 데 평균 66일이 걸린다고 보고했습니다. (개인/난이도별 18~254일 소요 편차 존재)", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf2, "• 안정적 항상성 형성 구간", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf2, "훈련을 무의식적 일상 루틴으로 정착시켜, 좋아졌다 무너지는 일시적 적응의 한계를 넘는 전략적 기간으로 활용합니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3단계: 100일 수초화기
    create_card_shape(s8, Inches(8.8), Inches(1.5), Inches(3.7), card_h)
    tb3 = s8.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "🌳 100일 뇌 재배선 (Myelination)"
    tf3.paragraphs[0].font.name = FONT_NAME
    tf3.paragraphs[0].font.size = Pt(15)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_ACCENT
    tf3.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf3, "• 신경가소성 및 수초화 원리 참고", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf3, "반복적인 자극은 신경의 정보 처리 속도를 가속하는 축삭 수초화(Myelination)를 유도하여 신경 연결을 강화하는 데 기여합니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf3, "• 100일의 실천적 목표 설계", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf3, "뉴런 결속을 유도하는 헤브의 법칙(Hebbian learning)을 반영하여, 삶의 근본적인 인지 필터와 태도를 리모델링하는 실천적 기간입니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 9: 뇌교육 명상(BEM)의 임상적 실효성 및 유전자 바이오마커
    # ==================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, COLOR_LIGHT_BG)
    add_slide_header(s9, "뇌기반 명상(BEM) 프로그램의 인지·학습 및 뇌기능 연구 성과", "Educational & Brain Research")

    # 좌측: 인지·자기조절 자기주도학습 연구
    create_card_shape(s9, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_cl = s9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_cl = tb_cl.text_frame
    tf_cl.word_wrap = True
    tf_cl.paragraphs[0].text = "🎓 자기주도학습력 및 인지 조절 향상"
    tf_cl.paragraphs[0].font.name = FONT_NAME
    tf_cl.paragraphs[0].font.size = Pt(16)
    tf_cl.paragraphs[0].font.bold = True
    tf_cl.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_cl.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_cl, "• 국내 뇌교육 프로그램 아동/청소년 연구", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_cl, "중학생 등을 대상으로 8주간의 뇌교육 프로그램을 진행한 결과, 인지 및 감정조절을 바탕으로 한 '자기주도학습능력'이 유의하게 증가된 성과가 보고되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_cl, "• 뇌기반 교육의 실천적 타당성 확보", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_cl, "신경과학, 인지심리, 교육심리학 등의 연구 성과를 현장에 연결하여 학습자의 메타인지와 자기 조절력을 키우는 실천 모델로 입증되었습니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # 우측: 생리적·뇌 기능 변화 연구 (명상 리포트)
    create_card_shape(s9, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_cr = s9.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_cr = tb_cr.text_frame
    tf_cr.word_wrap = True
    tf_cr.paragraphs[0].text = "🧠 생리학적 지표 및 뇌 기능 활성화"
    tf_cr.paragraphs[0].font.name = FONT_NAME
    tf_cr.paragraphs[0].font.size = Pt(16)
    tf_cr.paragraphs[0].font.bold = True
    tf_cr.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_cr.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_cr, "• 뇌교육 기반 명상 연구 리포트 종합", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_cr, "한국식 명상 관련 학술 연구들을 분석한 자료에 의하면, 수련 전후 스트레스 척도 완화, 긍정적 정서의 유의미한 향상이 지속적으로 보고되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_cr, "• 도파민 및 디폴트 모드 네트워크 변화", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_cr, "혈장 내 도파민 활성화와 뇌 디폴트 모드 네트워크(DMN)의 기능적 연결성 변화 등, 생리 및 뇌 신호 조절을 시사하는 유의미한 축적형 데이터가 축적되어 있습니다.", 11, False, COLOR_TEXT_MUTED, 10)

    # 하단 각주 (주의사항) 추가
    add_paragraph_to_tf(tf_cr, "※ 만성 대사 및 특정 염증 유전자 발현 지표 등의 수치 데이터는 내부 예비 연구 결과로서, 대외 제출용 공식 보고서 작성 시에는 연구 설계를 병기합니다.", 8.5, False, COLOR_TEXT_MUTED, 0)

    # ==================================================
    # SLIDE 10: 대규모 성과 및 공공/기업 도입 레퍼런스
    # ==================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, COLOR_LIGHT_BG)
    add_slide_header(s10, "대규모 점검 성과 및 신뢰도 검증 레퍼런스", "References & Track Record")

    # 1. 성과 요약 카드 (좌측)
    create_card_shape(s10, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.1))
    tb_stats = s10.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(4.1), Inches(4.7))
    tf_stats = tb_stats.text_frame
    tf_stats.word_wrap = True
    tf_stats.paragraphs[0].text = "🚀 단기 대규모 측정 성과"
    tf_stats.paragraphs[0].font.name = FONT_NAME
    tf_stats.paragraphs[0].font.size = Pt(16)
    tf_stats.paragraphs[0].font.bold = True
    tf_stats.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_stats.paragraphs[0].space_after = Pt(20)

    add_paragraph_to_tf(tf_stats, "출시 2달 미만 누적 측정자", 12, True, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_stats, "8,000명 돌파", 32, True, COLOR_ACCENT, 16)
    
    add_paragraph_to_tf(tf_stats, "“단 2개월 만에 8,000명 이상의 고유 점검 데이터를 안정적으로 누적했습니다. 이는 하드웨어 리소스 최적화 및 2D 컴퓨터 비전 실시간 연산 안전성이 실전 현장에서 완벽하게 검증되었음을 대변합니다.”", 11, False, COLOR_TEXT_DARK, 6)

    # 2. 공공기관 및 기업 레퍼런스 카드 (우측)
    create_card_shape(s10, Inches(5.6), Inches(1.5), Inches(6.9), Inches(5.1))
    tb_ref = s10.shapes.add_textbox(Inches(5.8), Inches(1.7), Inches(6.5), Inches(4.7))
    tf_ref = tb_ref.text_frame
    tf_ref.word_wrap = True
    tf_ref.paragraphs[0].text = "🏢 관공서 및 대기업·단체 도입 레퍼런스"
    tf_ref.paragraphs[0].font.name = FONT_NAME
    tf_ref.paragraphs[0].font.size = Pt(16)
    tf_ref.paragraphs[0].font.bold = True
    tf_ref.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_ref.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_ref, "• 주요 지자체 복지 사업 지원", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "춘천시청, 포항시청, 춘천교육청, 천안시 행복키움단 등 시청 및 교육 기관 산하 건강/복지 프로그램으로 도입 및 운영되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_ref, "• 전문 의료진 및 체육·공익 단체", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "서울약사협회, 대구약사협회, 경북약사협회, 평택체육회, 안동노인회 등 높은 객관성이 요구되는 약사 및 보건 체육 단체 회원들을 대상으로 성공적인 단체 건강 지표 수립 검사를 실행했습니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_ref, "• 대기업 복지 연계", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "IT 기업인 카카오(Kakao) 등 소속 임직원의 멘탈 웰니스 측정 및 직장인 증후군 자가 예방 지원 툴로써 AI 코드맵 점검이 널리 기여하고 있습니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # PPT 파일 저장
    outputs = ['부경대학교_3바디7코드_AI코드맵_신뢰성보고서.pptx', '부경대학교_3바디7코드_AI코드맵_신뢰성보고서_v2.pptx']
    for output_ppt in outputs:
        try:
            prs.save(output_ppt)
            print(f"[완료] 부경대학교 보고서 PPTX 생성 성공: {output_ppt}")
        except PermissionError:
            print(f"[알림] {output_ppt} 파일이 현재 사용 중이어서 덮어쓰기를 건너뜁니다.")
    print(f"\n==================================================")
    print(f"==================================================")

if __name__ == '__main__':
    generate_pukyong_presentation()
