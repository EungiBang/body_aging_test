# Excel 통계 데이터를 기반으로 세련된 레이아웃의 PPTX 보고서를 자동 생성하는 스크립트
import json
import os
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 테마 색상 정의 (세련된 비즈니스 딥 네이비 테마)
COLOR_DARK_BG = RGBColor(15, 23, 42)      # 슬레이트 블랙/딥 네이비
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # 은은한 그레이시 화이트
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PRIMARY = RGBColor(31, 73, 125)     # 메인 네이비 블루
COLOR_SECONDARY = RGBColor(59, 130, 246)  # 서브 스카이 블루
COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # 차콜
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 뮤트 그레이
COLOR_ACCENT = RGBColor(245, 158, 11)     # 오렌지 액센트
COLOR_CARD_BORDER = RGBColor(226, 232, 240)

FONT_NAME = "Malgun Gothic"

def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title, subtitle):
    # 상단 대분류/서브 타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.3))
    p_sub = sub_box.text_frame.paragraphs[0]
    p_sub.text = subtitle.upper()
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(9.5)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.6))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_DARK

def create_card_shape(slide, left, top, width, height, bg_color=COLOR_WHITE):
    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = bg_color
    card_shape.line.color.rgb = COLOR_CARD_BORDER
    card_shape.line.width = Pt(1)
    return card_shape

def format_cell_content(cell, text, font_size=10, is_bold=False, font_color=COLOR_TEXT_DARK, alignment=PP_ALIGN.CENTER):
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color

def style_table_header(table, col_count):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

def load_data_from_excel():
    excel_file = 'BTC_코드맵_점검_통계보고서_수정.xlsx'
    if not os.path.exists(excel_file):
        print(f"[오류] {excel_file} 파일이 없습니다.")
        return None

    wb = openpyxl.load_workbook(excel_file, data_only=True)
    ws = wb['지역_지점별_통계']

    raw_data = []
    current_region = "미지정"

    for r in range(2, ws.max_row + 1):
        val_region = ws.cell(row=r, column=1).value
        val_branch = ws.cell(row=r, column=2).value
        val_raw = ws.cell(row=r, column=3).value
        val_uniq = ws.cell(row=r, column=4).value

        # 전체 합계행 도달 시 정지
        if val_region == '전체 합계' or val_branch == '전체 합계' or (val_region is not None and '합계' in str(val_region)):
            break

        if val_region is not None and str(val_region).strip() != '':
            current_region = str(val_region).strip()

        if val_branch:
            raw_data.append({
                'region': current_region,
                'branch': str(val_branch).strip(),
                'raw_count': val_raw if isinstance(val_raw, int) else 0,
                'unique_count': val_uniq if isinstance(val_uniq, int) else 0
            })

    print(f"로드된 지점별 데이터 수: {len(raw_data)}건")
    return raw_data

def generate_presentation():
    data = load_data_from_excel()
    if not data:
        return

    # 통계 요약 집계
    total_branches = len(data)
    total_raw_count = sum(item['raw_count'] for item in data)
    total_uniq_count = sum(item['unique_count'] for item in data)
    
    # 지역별 집계
    region_stats = {}
    for item in data:
        r = item['region']
        if r not in region_stats:
            region_stats[r] = {'raw': 0, 'uniq': 0, 'branches': []}
        region_stats[r]['raw'] += item['raw_count']
        region_stats[r]['uniq'] += item['unique_count']
        region_stats[r]['branches'].append(item)

    total_regions = len(region_stats)

    # 지점별 점검인원 순위 (Top 5)
    sorted_branches = sorted(data, key=lambda x: x['unique_count'], reverse=True)
    top_branches = sorted_branches[:5]

    # PPT 생성 시작
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==================================================
    # SLIDE 1: 표지 슬라이드
    # ==================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_DARK_BG)

    # 포인트 좌측 바
    left_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_PRIMARY
    left_bar.line.fill.background()

    # 상단 장식 데코레이션 박스
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(1.8), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_SECONDARY
    badge.line.fill.background()
    badge.text_frame.text = "점검 결과 보고"
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = FONT_NAME
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_WHITE

    # 제목 텍스트 상자
    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "CODEMAP AI 점검 데이터 정밀 분석"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(16)
    p0.font.color.rgb = COLOR_SECONDARY
    p0.font.bold = True
    p0.space_after = Pt(12)

    p1 = tf.add_paragraph()
    p1.text = "지점별 코드맵 점검 통계 보고서"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "중복 제외 실제 점검 인원 기준 분석"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_WHITE
    p2.font.bold = False
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = f"분석 대상: 24개 지역 / 139개 지점 (총 {total_raw_count:,}건 점검 / 고유 사용자 {total_uniq_count:,}명)"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(6)

    p4 = tf.add_paragraph()
    p4.text = "기준 일자: 2026. 06. 23.  |  데이터 소스: BTC 중앙관리서버"
    p4.font.name = FONT_NAME
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_TEXT_MUTED

    # ==================================================
    # SLIDE 2: 전국 점검 통계 요약
    # ==================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_LIGHT_BG)
    add_slide_header(s2, "전국 코드맵 점검 현황 요약", "Overall Summary")

    # 카드 1: 점검 규모 요약
    create_card_shape(s2, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.1))
    tb_c1 = s2.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(3.0), Inches(4.5))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    
    p = tf_c1.paragraphs[0]
    p.text = "📊 누적 점검 통계"
    p.font.name = FONT_NAME
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(20)

    for label, val, is_accent in [
        ("총 고유 인원수", f"{total_uniq_count:,} 명", True),
        ("총 점검 건수", f"{total_raw_count:,} 건", False),
        ("중복 점검 건수", f"{total_raw_count - total_uniq_count:,} 건", False),
        ("중복 점검 비율", f"{round((1 - total_uniq_count/total_raw_count)*100, 1)}%", False)
    ]:
        p_lbl = tf_c1.add_paragraph()
        p_lbl.text = label
        p_lbl.font.name = FONT_NAME
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = COLOR_TEXT_MUTED
        p_lbl.space_after = Pt(2)

        p_val = tf_c1.add_paragraph()
        p_val.text = val
        p_val.font.name = FONT_NAME
        p_val.font.size = Pt(24) if is_accent else Pt(16)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_ACCENT if is_accent else COLOR_TEXT_DARK
        p_val.space_after = Pt(12)

    # 카드 2: 점검 인프라 요약
    create_card_shape(s2, Inches(4.8), Inches(1.5), Inches(3.6), Inches(5.1))
    tb_c2 = s2.shapes.add_textbox(Inches(5.1), Inches(1.8), Inches(3.0), Inches(4.5))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True

    p = tf_c2.paragraphs[0]
    p.text = "📍 점검 지점 분포"
    p.font.name = FONT_NAME
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(20)

    for label, val in [
        ("활성화 지역 수", f"{total_regions:,} 개 지역"),
        ("총 설치 지점 수", f"{total_branches:,} 개 지점"),
        ("지점 평균 점검수", f"{round(total_uniq_count/total_branches, 1)} 명 / 지점"),
    ]:
        p_lbl = tf_c2.add_paragraph()
        p_lbl.text = label
        p_lbl.font.name = FONT_NAME
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = COLOR_TEXT_MUTED
        p_lbl.space_after = Pt(2)

        p_val = tf_c2.add_paragraph()
        p_val.text = val
        p_val.font.name = FONT_NAME
        p_val.font.size = Pt(20)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_TEXT_DARK
        p_val.space_after = Pt(16)

    # 카드 3: 점검 TOP 5 지점
    create_card_shape(s2, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.1))
    tb_c3 = s2.shapes.add_textbox(Inches(9.1), Inches(1.8), Inches(3.1), Inches(4.5))
    tf_c3 = tb_c3.text_frame
    tf_c3.word_wrap = True

    p = tf_c3.paragraphs[0]
    p.text = "🏆 최다 점검 지점 TOP 5"
    p.font.name = FONT_NAME
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(18)

    for idx, b_data in enumerate(top_branches, 1):
        p_top = tf_c3.add_paragraph()
        p_top.text = f"{idx}위.  {b_data['branch']} 지점"
        p_top.font.name = FONT_NAME
        p_top.font.size = Pt(12)
        p_top.font.bold = True
        p_top.font.color.rgb = COLOR_TEXT_DARK

        p_top_sub = tf_c3.add_paragraph()
        p_top_sub.text = f"    지역: {b_data['region']}  |  점검인원: {b_data['unique_count']}명 (총 {b_data['raw_count']}건)"
        p_top_sub.font.name = FONT_NAME
        p_top_sub.font.size = Pt(10.5)
        p_top_sub.font.color.rgb = COLOR_TEXT_MUTED
        p_top_sub.space_after = Pt(10)

    # ==================================================
    # SLIDE 3: 지역별 점검 규모 비교
    # ==================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_LIGHT_BG)
    add_slide_header(s3, "전국 지역별 코드맵 점검 현황 비교", "Regional Ranking")

    # 지역별 고유 인원 기준 정렬
    sorted_regions = sorted(region_stats.items(), key=lambda x: x[1]['uniq'], reverse=True)
    
    # 24개 지역을 좌/우 2개 테이블(각 12개 행)로 나누어 한 슬라이드에 깔끔하게 나열
    rows = 13 # 헤더 1 + 데이터 12
    
    # 좌측 테이블 (1~12위)
    t_shape_l = s3.shapes.add_table(rows, 4, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.3))
    t_l = t_shape_l.table
    t_l.columns[0].width = Inches(0.8)  # 순위
    t_l.columns[1].width = Inches(1.8)  # 지역명
    t_l.columns[2].width = Inches(1.5)  # 총 건수
    t_l.columns[3].width = Inches(1.5)  # 고유 인원
    
    # 우측 테이블 (13~24위)
    t_shape_r = s3.shapes.add_table(rows, 4, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.3))
    t_r = t_shape_r.table
    t_r.columns[0].width = Inches(0.8)
    t_r.columns[1].width = Inches(1.8)
    t_r.columns[2].width = Inches(1.5)
    t_r.columns[3].width = Inches(1.5)

    headers = ["순위", "지역명", "총 점검수", "점검 인원"]
    
    for t in [t_l, t_r]:
        for i, h in enumerate(headers):
            format_cell_content(t.cell(0, i), h, 10.5, True, COLOR_WHITE)
        style_table_header(t, 4)

    # 데이터 바인딩
    for idx, (r_name, r_info) in enumerate(sorted_regions, 1):
        if idx <= 12:
            target_table = t_l
            row_num = idx
        else:
            target_table = t_r
            row_num = idx - 12

        is_top = idx <= 5
        text_color = COLOR_PRIMARY if is_top else COLOR_TEXT_DARK
        bold_flag = is_top

        format_cell_content(target_table.cell(row_num, 0), str(idx), 10, bold_flag, text_color)
        format_cell_content(target_table.cell(row_num, 1), r_name, 10, bold_flag, text_color, PP_ALIGN.LEFT)
        format_cell_content(target_table.cell(row_num, 2), f"{r_info['raw']:,}건", 10, False, COLOR_TEXT_DARK, PP_ALIGN.RIGHT)
        
        # 고유 점검인원 볼드 강조
        format_cell_content(target_table.cell(row_num, 3), f"{r_info['uniq']:,}명", 10, True, COLOR_ACCENT if is_top else COLOR_PRIMARY, PP_ALIGN.RIGHT)

    # ==================================================
    # SLIDE 4 ~ 8: 지역별 세부 지점 현황 (슬라이드 당 4~5개 지역씩 분할)
    # ==================================================
    # 24개 지역을 5개 슬라이드로 나누기 위한 매핑 정의
    region_groups = [
        {
            "title": "서울 및 경기 북부 권역 세부 현황",
            "regions": ["본사", "서울강북1", "서울강북2", "서울강북3", "서울강남1", "서울강남2"]
        },
        {
            "title": "경기 남부 및 인천 권역 세부 현황",
            "regions": ["경기북부", "경기남부1", "경기남부2", "경기남부3", "인천"]
        },
        {
            "title": "강원, 대전 및 충청 권역 세부 현황",
            "regions": ["강원", "대전", "충북", "충남"]
        },
        {
            "title": "호남 및 대구·경북 권역 세부 현황",
            "regions": ["광주전남", "전북", "대구", "경북"]
        },
        {
            "title": "부산·울산·경남 및 기타 권역 세부 현황",
            "regions": ["경남", "울산", "부산", "제주", "단무도"]
        }
    ]

    for g_idx, group in enumerate(region_groups, 4):
        s_g = prs.slides.add_slide(blank_layout)
        set_slide_background(s_g, COLOR_LIGHT_BG)
        add_slide_header(s_g, group['title'], f"Detailed Branch Stats (Part {g_idx-3})")

        g_regions = group['regions']
        reg_count = len(g_regions)

        # 4개 지역: 2x2 그리드
        # 5개 지역: 상단 3개, 하단 2개
        # 6개 지역: 2x3 그리드 (가로 3, 세로 2)
        
        for r_idx, r_name in enumerate(g_regions):
            r_info = region_stats.get(r_name, {'uniq': 0, 'raw': 0, 'branches': []})
            
            # 그리드 위치 좌표 계산
            if reg_count == 4:
                # 2x2 배정
                col = r_idx % 2
                row = r_idx // 2
                left = Inches(0.8 + col * 5.9)
                top = Inches(1.5 + row * 2.7)
                width = Inches(5.6)
                height = Inches(2.4)
            elif reg_count == 5:
                # 상단 3개, 하단 2개 (좌우 쏠림 방지 오프셋)
                if r_idx < 3:
                    left = Inches(0.8 + r_idx * 3.9)
                    top = Inches(1.5)
                else:
                    left = Inches(2.75 + (r_idx - 3) * 3.9)
                    top = Inches(4.2)
                width = Inches(3.6)
                height = Inches(2.4)
            else: # 6개 지역
                # 2x3 배정
                col = r_idx % 3
                row = r_idx // 3
                left = Inches(0.8 + col * 3.9)
                top = Inches(1.5 + row * 2.7)
                width = Inches(3.6)
                height = Inches(2.4)

            # 카드 셰이프 그리기
            create_card_shape(s_g, left, top, width, height)

            # 타이틀 영역 추가
            tb_title = s_g.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4))
            tf_t = tb_title.text_frame
            tf_t.word_wrap = True
            p_t = tf_t.paragraphs[0]
            p_t.text = f"📍 {r_name} ({r_info['uniq']}명)"
            p_t.font.name = FONT_NAME
            p_t.font.size = Pt(12)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_PRIMARY

            # 지점 리스트 텍스트 상자 추가
            tb_body = s_g.shapes.add_textbox(left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55))
            tf_b = tb_body.text_frame
            tf_b.word_wrap = True

            # 각 지점 리스트 작성 시 오버플로우 방지 처리
            # 지점 개수가 많으면 쉼표(,) 구분형 2줄 또는 컬럼 분할 느낌으로 압축
            b_list = sorted(r_info['branches'], key=lambda x: x['unique_count'], reverse=True)
            
            p_b = tf_b.paragraphs[0]
            p_b.font.name = FONT_NAME
            p_b.font.size = Pt(10)
            p_b.font.color.rgb = COLOR_TEXT_DARK

            if len(b_list) <= 6:
                # 6개 지점 이하는 한 줄씩 리스트업
                for idx_b, b in enumerate(b_list):
                    if idx_b == 0:
                        p_line = p_b
                    else:
                        p_line = tf_b.add_paragraph()
                    
                    p_line.text = f"• {b['branch']}: {b['unique_count']}명"
                    p_line.font.name = FONT_NAME
                    p_line.font.size = Pt(10)
                    p_line.font.color.rgb = COLOR_TEXT_DARK
                    p_line.space_after = Pt(2)
            else:
                # 지점이 7개 이상이면 한 줄에 2개 지점씩 표기하여 카드 밖 이탈 방지
                # 2열 컴팩트 배치
                lines = []
                for idx_b in range(0, len(b_list), 2):
                    b1 = b_list[idx_b]
                    part1 = f"{b1['branch']}({b1['unique_count']}명)"
                    
                    if idx_b + 1 < len(b_list):
                        b2 = b_list[idx_b + 1]
                        part2 = f"{b2['branch']}({b2['unique_count']}명)"
                        lines.append(f"• {part1.ljust(15)} • {part2}")
                    else:
                        lines.append(f"• {part1}")
                
                for idx_l, line_txt in enumerate(lines):
                    if idx_l == 0:
                        p_line = p_b
                    else:
                        p_line = tf_b.add_paragraph()
                    p_line.text = line_txt
                    p_line.font.name = FONT_NAME
                    p_line.font.size = Pt(9.5)
                    p_line.font.color.rgb = COLOR_TEXT_DARK
                    p_line.space_after = Pt(2)

    # PPT 파일 저장
    output_ppt = 'BTC_지점별_코드맵_점검현황.pptx'
    prs.save(output_ppt)
    print(f"\n==================================================")
    print(f"[완료] 파워포인트 보고서 생성 완료: {output_ppt}")
    print(f"==================================================")

if __name__ == '__main__':
    generate_presentation()
