# 학교 선생님 대상 3바디 7코드 AI 30분 강의용 PPTX 생성 스크립트 (11장 확장본)
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 테마 색상 (신뢰감을 주는 프리미엄 딥 네이비 & 아일랜드 블루 테마)
COLOR_DARK_BG = RGBColor(15, 23, 42)      # 슬레이트 블랙/딥 네이비
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # 은은한 그레이시 화이트
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PRIMARY = RGBColor(27, 54, 93)      # 다크 네이비
COLOR_SECONDARY = RGBColor(59, 130, 246)  # 아일랜드 블루
COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # 차콜 블랙
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 뮤트 그레이
COLOR_ACCENT = RGBColor(245, 158, 11)     # 골드/오렌지 액센트
COLOR_CARD_BORDER = RGBColor(226, 232, 240)
COLOR_LIGHT_BLUE = RGBColor(239, 246, 255)

FONT_NAME = "Malgun Gothic"

def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title, subtitle):
    # 상단 대분류/서브 타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.3))
    p_sub = sub_box.text_frame.paragraphs[0]
    p_sub.text = subtitle.upper()
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(9.5)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.6))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_DARK

def create_card_shape(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER):
    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = bg_color
    card_shape.line.color.rgb = border_color
    card_shape.line.width = Pt(1)
    return card_shape

def add_paragraph_to_tf(tf, text, font_size=11, is_bold=False, font_color=COLOR_TEXT_DARK, space_after=6, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    return p

def format_cell_content(cell, text, font_size=10, is_bold=False, font_color=COLOR_TEXT_DARK, alignment=PP_ALIGN.CENTER):
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = font_color

def style_table_header(table, col_count):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

def generate_lecture_presentation():
    # PPT 생성 시작 (Widescreen 16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==================================================
    # SLIDE 1: 표지 슬라이드 (Dark BG)
    # ==================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_DARK_BG)

    left_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_SECONDARY
    left_bar.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(2.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_SECONDARY
    badge.line.fill.background()
    badge.text_frame.text = "교사 홀리스틱 웰니스"
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = FONT_NAME
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_WHITE

    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "3바디 7코드 AI로 만나는 새로운 건강 패러다임"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(16)
    p0.font.color.rgb = COLOR_SECONDARY
    p0.font.bold = True
    p0.space_after = Pt(12)

    p1 = tf.add_paragraph()
    p1.text = "교직 생활의 에너지를 깨우는 30분 웰니스 가이드"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "몸·마음·뇌의 통합적 균형과 7코드 에너지 기반 자가 컨디션 관리 솔루션"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_WHITE
    p2.font.bold = False
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "공동 연구 개발: 국가공인 브레인트레이너 전문가 단체 · 한국뇌과학연구원 · 글로벌사이버대학교"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(6)

    p4 = tf.add_paragraph()
    p4.text = "제공 기관: 브레인트레이닝센터"
    p4.font.name = FONT_NAME
    p4.font.size = Pt(12)
    p4.font.color.rgb = COLOR_TEXT_MUTED

    # ==================================================
    # SLIDE 2: 왜 몸·마음·뇌의 통합적 건강인가? (Light BG)
    # ==================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_LIGHT_BG)
    add_slide_header(s2, "신체·마음·뇌의 통합적 균형: 홀리스틱 웰니스", "Integrated Health Needs")

    card_w = Inches(3.6)
    card_h = Inches(5.1)

    # 1. 몸 (신체 건강에 치중된 관리의 한계)
    create_card_shape(s2, Inches(0.8), Inches(1.5), card_w, card_h)
    tb_b1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b1 = tb_b1.text_frame
    tf_b1.word_wrap = True
    tf_b1.paragraphs[0].text = "💪 육체 중심 관리의 한계"
    tf_b1.paragraphs[0].font.name = FONT_NAME
    tf_b1.paragraphs[0].font.size = Pt(15)
    tf_b1.paragraphs[0].font.bold = True
    tf_b1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b1.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b1, "• 반쪽짜리 건강 관리", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b1, "요가, 헬스 등 신체 단련만으로는 쌓인 스트레스와 피로를 온전히 해결할 수 없습니다. 몸의 피로는 정신적 피로와 유기적으로 연결되어 있습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b1, "• 몸과 뇌의 피드백 작용", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b1, "육체의 만성적 긴장은 뇌의 자율신경 조절 기능을 떨어뜨려 스트레스 민감도를 높집니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2. 마음 (교사 스트레스와 번아웃 실태)
    create_card_shape(s2, Inches(4.8), Inches(1.5), card_w, card_h)
    tb_b2 = s2.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_b2 = tb_b2.text_frame
    tf_b2.word_wrap = True
    tf_b2.paragraphs[0].text = "❤️ 교직 생활 속 번아웃과 스트레스"
    tf_b2.paragraphs[0].font.name = FONT_NAME
    tf_b2.paragraphs[0].font.size = Pt(15)
    tf_b2.paragraphs[0].font.bold = True
    tf_b2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b2.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b2, "• 일상 스트레스 인지율 25.9%", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b2, "질병관리청 통계에 따르면 성인 4명 중 1명은 일상 중 높은 수준의 스트레스를 경험합니다. 교사는 학생 지도와 행정 업무로 인해 그 노출 빈도가 더 높습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b2, "• 교사 번아웃 경험률 60~70% 돌파", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b2, "감정 노동과 직무 과부하로 정서적 고갈과 무기력을 느끼는 교직원이 매년 빠르게 증가하고 있습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3. 뇌 (현대 의과학적 입증 근거)
    create_card_shape(s2, Inches(8.8), Inches(1.5), Inches(3.7), card_h)
    tb_b3 = s2.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf_b3 = tb_b3.text_frame
    tf_b3.word_wrap = True
    tf_b3.paragraphs[0].text = "🧠 현대 의과학적 연계성"
    tf_b3.paragraphs[0].font.name = FONT_NAME
    tf_b3.paragraphs[0].font.size = Pt(15)
    tf_b3.paragraphs[0].font.bold = True
    tf_b3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_b3.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_b3, "• 심신신경면역학 (PNI) 입증", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b3, "스트레스 자극은 뇌의 HPA 축을 통해 코르티솔 분비를 유도하고, 신체의 면역 저하 및 소화기 장애로 직결됩니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_b3, "• Nature 2023 SCAN 연구", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_b3, "대뇌 1차 운동 피질 내부에는 신체 움직임, 인지 제어, 그리고 내부 생리 상태(혈압 등)를 동시 조절하는 기능적 네트워크의 존재가 보고되었습니다.", 10.5, False, COLOR_TEXT_MUTED, 8)

    # ==================================================
    # SLIDE 3: 오늘 나의 상태를 지배하는 '컨디션과 에너지' (Light BG)
    # ==================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_LIGHT_BG)
    add_slide_header(s3, "컨디션을 결정하는 열쇠: 생체 에너지와 7코드", "Energy & 7Code Concept")

    # 1. 생체 에너지 흐름 (좌측 카드)
    create_card_shape(s3, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_el = s3.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_el = tb_el.text_frame
    tf_el.word_wrap = True
    tf_el.paragraphs[0].text = "🔄 컨디션을 결정하는 에너지 순환의 고리"
    tf_el.paragraphs[0].font.name = FONT_NAME
    tf_el.paragraphs[0].font.size = Pt(16)
    tf_el.paragraphs[0].font.bold = True
    tf_el.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_el.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_el, "• 에너지가 저하될 때 (악순환)", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_el, "에너지가 떨어지면 기분 기복이 심해지고, 이는 자율신경계를 자극해 소화장애, 근육 긴장을 유발합니다. 이는 결국 뇌의 스트레스 인지 민감도를 증대시킵니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    
    add_paragraph_to_tf(tf_el, "• 에너지가 활성화될 때 (선순환)", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_el, "체내 에너지 흐름이 활발해지면 정서가 이완되고, 신체 대사 기능과 자가 치유 능력이 극대화됩니다. 뇌파가 가라앉으며 의사결정 속도와 활력이 향상됩니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2. 동양의 지혜와 7코드가 만나는 지점 (우측 카드)
    create_card_shape(s3, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_er = s3.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_er = tb_er.text_frame
    tf_er.word_wrap = True
    tf_er.paragraphs[0].text = "📜 고대의 신체 이론에서 과학적 7코드로"
    tf_er.paragraphs[0].font.name = FONT_NAME
    tf_er.paragraphs[0].font.size = Pt(16)
    tf_er.paragraphs[0].font.bold = True
    tf_er.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_er.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_er, "• 동양 철학의 '3단전'과 인도의 '7차크라'", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_er, "고대의 전통 의학에서는 인체의 건강 중심점을 에너지가 모이는 세 군데의 '단전(하단전, 중단전, 상단전)'과 '일곱 개의 에너지 바퀴(차크라)'로 설명해왔습니다.", 10.5, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_er, "• 현대 해부학을 기반으로 계량화한 7CODE", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_er, "브레인트레이닝센터는 이 에너지 중심점들의 실체가 우리 몸의 '자율신경총(Nerve Plexus)' 및 '호르몬 내분비선'의 위치와 일치함에 주목하여, 일반인 누구나 자신의 에너지 밸런스를 직관적으로 확인할 수 있도록 계량화했습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 4: 7CODE 에너지 시스템의 해부학적 매핑 (Light BG)
    # ==================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, COLOR_LIGHT_BG)
    add_slide_header(s4, "7CODE 에너지 시스템의 자율신경총 및 내분비선 해석 구조", "7Code Anatomical Mapping")

    # 설명 텍스트
    desc_box = s4.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.6))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.paragraphs[0].text = "7CODE 시스템은 인체의 주요 신경 다발과 호르몬 분비관을 7단계 에너지 센터로 규정하여, 각 영역이 지배하는 생리적 및 정서적 밸런스를 측정하는 도구입니다."
    tf_desc.paragraphs[0].font.name = FONT_NAME
    tf_desc.paragraphs[0].font.size = Pt(11)
    tf_desc.paragraphs[0].font.bold = False
    tf_desc.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED

    # 테이블 생성 (8행 5열)
    t_shape = s4.shapes.add_table(8, 5, Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.9))
    table = t_shape.table
    table.columns[0].width = Inches(1.6)  # 코드 구분
    table.columns[1].width = Inches(2.4)  # 자율신경총
    table.columns[2].width = Inches(2.2)  # 주요 내분비선
    table.columns[3].width = Inches(4.9)  # 주요 기능 (생리적/심리적)
    table.columns[4].width = Inches(1.0)  # 기준

    headers = ["코드 구분", "대응 자율신경총", "대응 내분비선", "주요 기능 (생리적/심리적)", "설명모델"]
    for i, h in enumerate(headers):
        format_cell_content(table.cell(0, i), h, 11, True, COLOR_WHITE)
    style_table_header(table, 5)

    data_7code = [
        ("1CODE (기초/생존)", "골반 신경총 (Pelvic Plexus)", "부신 (Adrenal Gland)", "스트레스 대응, 아드레날린 분비 조절, 하체 밸런스 및 생존 활력", "생리학적"),
        ("2CODE (감정/관계)", "하복부/요천추 신경총", "생식선 (Ovaries/Testes)", "에너지 및 체액의 하체 순환, 타인 수용력, 감정 기복 조절", "상관성을"),
        ("3CODE (의욕/추진)", "복강/태양신경총", "췌장 (Pancreas)", "대사 작용 및 소화기능 지배, 주체적인 통제력, 자존감", "이해하기"),
        ("4CODE (사랑/안정)", "심장 신경총 (Cardiac Plexus)", "흉선 (Thymus Gland)", "면역 활성 조절 지원, 가슴 답답함 완화, 정서적 회복탄력성", "위한"),
        ("5CODE (소통/표현)", "경부/인두 신경총", "갑상선 (Thyroid Gland)", "기초대사량 조절, 창의적 자기 표현, 목과 어깨 긴장 완화", "통합적"),
        ("6CODE (집중/통찰)", "두개강 내 신경망 (송과체 주변)", "뇌하수체 (Pituitary Gland)", "수면 주기 및 호르몬 총괄 통제, 감정 통제 및 판단력", "참고적"),
        ("7CODE (방향/의미)", "대뇌피질 및 중추신경계", "송과체 (Pineal Gland)", "고차원 인지 기능, 삶의 가치와 정체성 방향성 수립", "설명모델")
    ]

    for idx, row_data in enumerate(data_7code, 1):
        bg_col = COLOR_WHITE if idx % 2 != 0 else COLOR_LIGHT_BLUE
        for col_idx, text in enumerate(row_data):
            cell = table.cell(idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_col
            
            align = PP_ALIGN.CENTER if col_idx in [0, 1, 2, 4] else PP_ALIGN.LEFT
            is_bold = col_idx == 0
            text_color = COLOR_PRIMARY if col_idx == 0 else COLOR_TEXT_DARK
            format_cell_content(cell, text, 10, is_bold, text_color, align)

    # ==================================================
    # SLIDE 5: 3바디 7코드 AI 정밀 측정 핵심 기술 (Light BG)
    # ==================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, COLOR_LIGHT_BG)
    add_slide_header(s5, "카메라 하나로 파악하는 AI 컴퓨터 비전 정밀 트래킹", "AI Technology & Accuracy")

    # 상단 요약 설명 박스
    create_card_shape(s5, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1), COLOR_LIGHT_BLUE, COLOR_SECONDARY)
    tb_top = s5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.9))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = True
    p_top = tf_top.paragraphs[0]
    p_top.text = "💡 독자적 AI 기술 탑재로 측정의 신뢰성과 객관성 확보 (관련 특허출원 신청 진행 중)"
    p_top.font.name = FONT_NAME
    p_top.font.size = Pt(13)
    p_top.font.bold = True
    p_top.font.color.rgb = COLOR_PRIMARY
    p_top.space_after = Pt(4)
    add_paragraph_to_tf(tf_top, "모바일/태블릿 기기의 웹 환경에서 별도의 외부 센서 부착 없이, 오직 카메라를 기반으로 한 실시간 물리 데이터 추출 엔진을 탑재해 3바디 상태를 정확하게 추적합니다.", 10.5, False, COLOR_TEXT_DARK, 0)

    # 하단 3단 상세 기술 카드
    card_y = Inches(2.7)
    card_h2 = Inches(4.0)
    card_w2 = Inches(3.6)

    # 기술 1: 실시간 부정방지(Cheat Prevention) 기술
    create_card_shape(s5, Inches(0.8), card_y, card_w2, card_h2)
    tb1 = s5.shapes.add_textbox(Inches(1.0), card_y + Inches(0.2), Inches(3.2), card_h2 - Inches(0.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "🛡️ 실시간 움직임 흔들림 감지"
    tf1.paragraphs[0].font.name = FONT_NAME
    tf1.paragraphs[0].font.size = Pt(14)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf1.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf1, "• 미세 흔들림 추적 (Sway Score)", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf1, "자세를 유지할 때 발생하는 미세한 좌표 흔들림을 추적하여 흔들림 계수(Sway Score)로 연산해 밸런스를 측정합니다.", 10, False, COLOR_TEXT_MUTED, 8)
    add_paragraph_to_tf(tf1, "• Foot Drop 센싱 기술 적용", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf1, "양 발목 Y축 좌표의 실시간 프레임 변동을 정밀 검지하여 외발 서기 중 발이 닿는 행위를 자동 판정합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 기술 2: WebGL 가속 및 오클루전 극복
    create_card_shape(s5, Inches(4.8), card_y, card_w2, card_h2)
    tb2 = s5.shapes.add_textbox(Inches(5.0), card_y + Inches(0.2), Inches(3.2), card_h2 - Inches(0.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "⚡ 최적화 연산 및 가려짐 극복"
    tf2.paragraphs[0].font.name = FONT_NAME
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf2.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf2, "• WebGL 메모리 관리 아키텍처", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf2, "매 프레임 생성되는 영상 텍스처 데이터를 실시간으로 소멸시켜 저사양 웹 브라우저에서도 끊김 현상이 발생하지 않습니다.", 10, False, COLOR_TEXT_MUTED, 8)
    add_paragraph_to_tf(tf2, "• 관절 가려짐(Occlusion) 보정 알고리즘", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf2, "운동 수행 중 관절이 겹치는 부위를 수학적 최적 추정 기법으로 보완하여 끊김 없는 좌표 추출이 가능합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # 기술 3: 가중치 기반 신체나이 환산 알고리즘
    create_card_shape(s5, Inches(8.8), card_y, Inches(3.7), card_h2)
    tb3 = s5.shapes.add_textbox(Inches(9.0), card_y + Inches(0.2), Inches(3.3), card_h2 - Inches(0.4))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "📊 신체/체력 나이 역산 시스템"
    tf3.paragraphs[0].font.name = FONT_NAME
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf3.paragraphs[0].space_after = Pt(14)
    add_paragraph_to_tf(tf3, "• 연령/성별 대조 정규화 스코어링", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf3, "한국인의 성별, 연령대별 평균 달성 표준에 따라 선형 보간 채점하여 관대한 임의 채점을 원천 차단합니다.", 10, False, COLOR_TEXT_MUTED, 8)
    add_paragraph_to_tf(tf3, "• 다차원 가중치 결합 공식 적용", 11, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf3, "근력, 유연성, 자세 제어, 밸런스 등 여러 테스트의 달성 결과를 적정 비율의 가중치로 합산하여 최종 건강 연령을 도출합니다.", 10, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 6: [NEW] 눈감고 한발서기: 3BODY 건강 밸런스 측정 (Light BG)
    # ==================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, COLOR_LIGHT_BG)
    add_slide_header(s6, "눈감고 한발서기: 신체·마음·뇌의 종합 밸런스 테스트", "One-Leg Stand & 3Body Balance")

    card_w = Inches(3.6)
    card_h = Inches(5.1)

    # 1. 몸 (신체적 측면)
    create_card_shape(s6, Inches(0.8), Inches(1.5), card_w, card_h)
    tb_o1 = s6.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_o1 = tb_o1.text_frame
    tf_o1.word_wrap = True
    tf_o1.paragraphs[0].text = "💪 몸(신체): 하체 근력과 정렬"
    tf_o1.paragraphs[0].font.name = FONT_NAME
    tf_o1.paragraphs[0].font.size = Pt(15)
    tf_o1.paragraphs[0].font.bold = True
    tf_o1.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_o1.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_o1, "• 건강나이의 직접적 지표", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_o1, "하체 근력과 척추의 뼈대 정렬을 측정하는 대표적인 국민 체력 표준입니다. 하체 근력이 탄탄할수록 서 있을 때 든든한 안정감이 확보됩니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_o1, "• 기둥이 무너질 때의 흔들림", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_o1, "골반과 코어 근육이 약하면 발이 지면에 닿기 전에 신체가 크게 요동치게 됩니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2. 뇌 (신경학적 측면)
    create_card_shape(s6, Inches(4.8), Inches(1.5), card_w, card_h, COLOR_WHITE, COLOR_SECONDARY)
    tb_o2 = s6.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf_o2 = tb_o2.text_frame
    tf_o2.word_wrap = True
    tf_o2.paragraphs[0].text = "🧠 뇌(두뇌): 소뇌 균형감각"
    tf_o2.paragraphs[0].font.name = FONT_NAME
    tf_o2.paragraphs[0].font.size = Pt(15)
    tf_o2.paragraphs[0].font.bold = True
    tf_o2.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    tf_o2.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_o2, "• 시각 차단과 평형 유지 능력", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_o2, "눈을 감는 순간 외부 시각 정보가 100% 차단되므로, 우리 뇌는 몸 전체 관절의 위치 감각(고유수용기 감각)에만 전적으로 의존하게 됩니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_o2, "• 소뇌의 실시간 평형 연산", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_o2, "소뇌(Cerebellum)는 내이의 전정감각과 척수의 운동 정보를 고속으로 연산하여 신체 균형을 실시간으로 바로잡습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3. 마음 (정서적/에너지 측면)
    create_card_shape(s6, Inches(8.8), Inches(1.5), Inches(3.7), card_h)
    tb_o3 = s6.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf_o3 = tb_o3.text_frame
    tf_o3.word_wrap = True
    tf_o3.paragraphs[0].text = "❤️ 마음(정서): 스트레스와 잡념"
    tf_o3.paragraphs[0].font.name = FONT_NAME
    tf_o3.paragraphs[0].font.size = Pt(15)
    tf_o3.paragraphs[0].font.bold = True
    tf_o3.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_o3.paragraphs[0].space_after = Pt(18)
    add_paragraph_to_tf(tf_o3, "• 불안과 뇌 잡념(DMN)의 방해", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_o3, "머릿속에 불필요한 생각과 스트레스가 많으면 뇌파가 요동쳐 순간적인 평정심과 집중력을 완전히 상실하게 됩니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_o3, "• 방전 시 15초 유지의 한계", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_o3, "컨디션(에너지)이 심각하게 저하되어 중심이 무너져 있는 상태라면 단 15초조차 눈을 감고 흔들림 없이 버티는 것은 대단히 어렵습니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 7: 스트레스를 조절하는 뇌기반 명상(BEM)의 기전 (Light BG)
    # ==================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, COLOR_LIGHT_BG)
    add_slide_header(s7, "뇌기반 명상(BEM)의 임상적 실효성 및 인지 조절 능력", "BEM Clinical Outcomes")

    # 좌측 카드: 정서조절 및 학업성취능력 (청소년 연구 사례)
    create_card_shape(s7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_l = s7.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "🎓 자기 조절력 및 주도성 향상 효과"
    tf_l.paragraphs[0].font.name = FONT_NAME
    tf_l.paragraphs[0].font.size = Pt(16)
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_l.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_l, "• 국내 아동/청소년 명상 수련 그룹 임상 연구", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l, "중학생을 대상으로 8주간 지속적인 뇌기반 명상을 지도한 결과, 기분 변화 제어력이 향상되어 '자기주도학습능력'이 유의하게 늘어났음이 확인되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_l, "• 메타인지 및 주의 제어력 발달", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_l, "자신의 감정 상태를 인지적으로 제3자 관점에서 내려다볼 수 있는 주의 제어 장치가 훈련을 통해 가시적으로 향상됩니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # 우측 카드: 뇌 기능 변화와 생리 지표 (성인 명상 연구 리포트)
    create_card_shape(s7, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    tb_r = s7.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "🧠 스트레스 이완 및 뇌 기능 활성화 기전"
    tf_r.paragraphs[0].font.name = FONT_NAME
    tf_r.paragraphs[0].font.size = Pt(16)
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_r.paragraphs[0].space_after = Pt(16)
    
    add_paragraph_to_tf(tf_r, "• DMN(디폴트 모드 네트워크) 과활성 억제", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r, "잡념과 불안이 가중될 때 흥분하는 뇌의 DMN 영역을 이완 훈련으로 제어하여, 쓸데없는 뇌의 피로와 과부하를 즉각적으로 리셋합니다.", 11, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf_r, "• 행복 전달 물질 도파민 및 호르몬 안정", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_r, "명상 수련 후 스트레스 척도 설문지 지표 개선은 물론 혈액 분석을 통해 긍정적 뇌 기능 분비물의 분비 향상이 입증되었습니다.", 11, False, COLOR_TEXT_MUTED, 10)
    add_paragraph_to_tf(tf_r, "※ 만성 대사 및 염증 발현 지표 수치는 내부 예비 연구 결과로서, 대외 발표시엔 연구 개요를 참조합니다.", 8.5, False, COLOR_TEXT_MUTED, 0)

    # ==================================================
    # SLIDE 8: 작심삼일을 극복하는 뇌과학적 3단계 습관 형성 (Light BG)
    # ==================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, COLOR_LIGHT_BG)
    add_slide_header(s8, "습관이 완전히 뇌에 정착되는 과정의 시간 기전", "Habit Formation & Brain Chemistry")

    card_w = Inches(3.6)
    card_h = Inches(5.1)

    # 1단계: 21일 적응기
    create_card_shape(s8, Inches(0.8), Inches(1.5), card_w, card_h)
    tb1 = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "🌱 1단계: 21일 적응기 (Adaptation)"
    tf1.paragraphs[0].font.name = FONT_NAME
    tf1.paragraphs[0].font.size = Pt(15)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    tf1.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf1, "• 자아 이미지 저항 수용", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf1, "맥스웰 몰츠 박사가 제안한 이론으로, 뇌가 낯선 행동이나 나의 변한 모습을 거부감 없이 인식하고 자아 이미지로 받아들이는 데 필요한 최초 적응 기간입니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf1, "• 뇌세포 생화학 신호 활성화", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf1, "명상과 트레이닝이 더 이상 어색하지 않은 수준의 시냅스 연결 초기 단계에 해당합니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 2단계: 66일 자동화기
    create_card_shape(s8, Inches(4.8), Inches(1.5), card_w, card_h, COLOR_WHITE, COLOR_SECONDARY)
    tb2 = s8.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "🌿 2단계: 66일 자동화기 (Automaticity)"
    tf2.paragraphs[0].font.name = FONT_NAME
    tf2.paragraphs[0].font.size = Pt(15)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf2.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf2, "• UCL 필리파 랠리 교수팀 입증", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf2, "특정 행동이 자발적 의식의 감시나 수고 없이 무의식적인 습관으로 자동 수행(Automaticity)되는 평균 도달 기간입니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf2, "• 향상된 기능의 항상성 유도", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf2, "하루 걸렀다고 즉각 무너지지 않고, 본래의 훈련 리듬으로 돌아오는 자발적인 조절 상태에 진입합니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # 3단계: 100일 수초화기
    create_card_shape(s8, Inches(8.8), Inches(1.5), Inches(3.7), card_h)
    tb3 = s8.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.7))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "🌳 3단계: 100일 수초화기 (Myelination)"
    tf3.paragraphs[0].font.name = FONT_NAME
    tf3.paragraphs[0].font.size = Pt(15)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_ACCENT
    tf3.paragraphs[0].space_after = Pt(18)
    
    add_paragraph_to_tf(tf3, "• 신경가소성에 따른 뇌 재배선 완료", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf3, "지속해서 자극을 받은 신경의 전달 전도 속도를 100배 가량 고속화시키는 '축삭 수초화(Myelination)'가 뇌에 구축되는 단계입니다.", 10.5, False, COLOR_TEXT_MUTED, 12)
    add_paragraph_to_tf(tf3, "• 뇌의 기본 활성 상태 리모델링", 11, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf3, "의도적으로 노력하지 않아도 항상 평정심과 컨디션을 훌륭하게 자가 정렬시키는 강력한 뇌의 결속 상태(헤브의 법칙)가 정착됩니다.", 10.5, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 9: 검증된 대규모 성과 및 공공/기업 도입 실적 (Light BG)
    # ==================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, COLOR_LIGHT_BG)
    add_slide_header(s9, "기관 도입의 타당성 및 신뢰도 검증 레퍼런스", "Public & Corporate References")

    # 1. 성과 요약 카드 (좌측)
    create_card_shape(s9, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.1))
    tb_stats = s9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(4.1), Inches(4.7))
    tf_stats = tb_stats.text_frame
    tf_stats.word_wrap = True
    tf_stats.paragraphs[0].text = "🚀 단기 대규모 검증 데이터 축적"
    tf_stats.paragraphs[0].font.name = FONT_NAME
    tf_stats.paragraphs[0].font.size = Pt(16)
    tf_stats.paragraphs[0].font.bold = True
    tf_stats.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_stats.paragraphs[0].space_after = Pt(20)

    add_paragraph_to_tf(tf_stats, "출시 2개월 만에 측정 누적 데이터", 12, True, COLOR_TEXT_MUTED, 4)
    add_paragraph_to_tf(tf_stats, "8,000명 돌파", 32, True, COLOR_ACCENT, 16)
    
    add_paragraph_to_tf(tf_stats, "“단기간에 수천 명의 고유 점검 결과가 안정적으로 적재된 것은, 카메라 2D 트래킹 기술과 Web 가속화 기술이 현장 실증 과정에서 완벽히 작동했음을 상징합니다.”", 11, False, COLOR_TEXT_DARK, 6)

    # 2. 공공기관 및 기업 레퍼런스 카드 (우측)
    create_card_shape(s9, Inches(5.6), Inches(1.5), Inches(6.9), Inches(5.1))
    tb_ref = s9.shapes.add_textbox(Inches(5.8), Inches(1.7), Inches(6.5), Inches(4.7))
    tf_ref = tb_ref.text_frame
    tf_ref.word_wrap = True
    tf_ref.paragraphs[0].text = "🏢 주요 도입 관공서 및 공익 단체 레퍼런스"
    tf_ref.paragraphs[0].font.name = FONT_NAME
    tf_ref.paragraphs[0].font.size = Pt(16)
    tf_ref.paragraphs[0].font.bold = True
    tf_ref.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_ref.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_ref, "• 공공 지방자치단체 및 교육 기관", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "춘천시청, 포항시청, 천안시 행복키움단, 그리고 춘천교육청의 복지 프로그램으로 도입되어, 지역 보건 및 교사 직무 연수 프로그램으로 성공리에 운용되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_ref, "• 전문 의료진 및 공익·보건 협회", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "서울약사협회, 대구약사협회, 경북약사협회 등 전문 약사 회원들을 비롯하여 평택체육회, 안동노인회 등 객관적 결과 지표를 요구하는 단체들에서 성공적으로 사용되었습니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_ref, "• IT 혁신 기업 복지 연계", 12, True, COLOR_TEXT_DARK, 6)
    add_paragraph_to_tf(tf_ref, "카카오(Kakao) 등 최첨단 IT 기업 소속 직장인들의 스트레스 증후군 조기 진단 및 멘탈 웰니스 강화의 실천 툴로 기여했습니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 10: 선생님을 위한 1분 뇌 휴식 및 에너지 리셋 팁 (Light BG)
    # ==================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, COLOR_LIGHT_BG)
    add_slide_header(s10, "일상에서 즉시 에너지를 조절하는 '1분 뇌 휴식법'", "1-Min Energy Reset Tip")

    card_w = Inches(5.6)
    card_h = Inches(5.1)

    # 좌측: 가슴 두드리기 기법 설명
    create_card_shape(s10, Inches(0.8), Inches(1.5), card_w, card_h)
    tb_sl = s10.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_sl = tb_sl.text_frame
    tf_sl.word_wrap = True
    tf_sl.paragraphs[0].text = "💓 가슴 두드리기 & '하' 방출 호흡법"
    tf_sl.paragraphs[0].font.name = FONT_NAME
    tf_sl.paragraphs[0].font.size = Pt(16)
    tf_sl.paragraphs[0].font.bold = True
    tf_sl.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    tf_sl.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_sl, "• 1단계: 가슴 중앙 4코드 정렬하기", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_sl, "허리를 곧게 펴고 편히 앉아, 가슴 흉선 부위인 4코드를 가벼운 주먹이나 손가락 끝으로 톡톡 부드럽게 두드려 줍니다. 내면의 무거운 뭉침 상태를 풀어냅니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_sl, "• 2단계: '하-' 방출구 호흡 실행", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_sl, "코로 맑은 산소를 가득 들이마신 후, 입을 가볍게 벌리고 '하-' 소리를 내며 몸 깊은 곳의 뜨거운 심화(心火) 스트레스 에너지를 길게 5회 토해 냅니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # 우측: 신경학적 원리
    create_card_shape(s10, Inches(6.9), Inches(1.5), card_w, card_h, COLOR_WHITE, COLOR_SECONDARY)
    tb_sr = s10.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.7))
    tf_sr = tb_sr.text_frame
    tf_sr.word_wrap = True
    tf_sr.paragraphs[0].text = "🧠 리셋 호흡법의 뇌과학적 원리"
    tf_sr.paragraphs[0].font.name = FONT_NAME
    tf_sr.paragraphs[0].font.size = Pt(16)
    tf_sr.paragraphs[0].font.bold = True
    tf_sr.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    tf_sr.paragraphs[0].space_after = Pt(16)

    add_paragraph_to_tf(tf_sr, "• 부교감 신경의 즉각적 활성화", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_sr, "가슴을 자극하고 긴 호흡을 내뱉는 행위는 교감 신경의 과부하로 뛰어오른 심장 박동을 낮추고 자율신경계 부교감 기능을 즉시 자극합니다.", 11, False, COLOR_TEXT_MUTED, 12)

    add_paragraph_to_tf(tf_sr, "• 뇌파를 안정 상태(알파파)로 정렬", 12, True, COLOR_TEXT_DARK, 4)
    add_paragraph_to_tf(tf_sr, "스트레스로 조급해져 베타파 상태로 치닫던 뇌 신경 활동이 단 1분의 집중과 이완 호흡을 통해 고요하고 평온한 '알파파(Alpha Waves)'로 전환됩니다.", 11, False, COLOR_TEXT_MUTED, 6)

    # ==================================================
    # SLIDE 11: 통합 건강(홀리스틱 웰니스)의 가치와 Q&A (Dark BG)
    # ==================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, COLOR_DARK_BG)

    left_bar = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_SECONDARY
    left_bar.line.fill.background()

    # 마무리 텍스트
    summary_box = s11.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(4.5))
    tf_sum = summary_box.text_frame
    tf_sum.word_wrap = True

    p0 = tf_sum.paragraphs[0]
    p0.text = "“선생님의 3BODY 균형이 교실 전체의 행복을 만듭니다”"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(22)
    p0.font.color.rgb = COLOR_SECONDARY
    p0.font.bold = True
    p0.space_after = Pt(20)

    p1 = tf_sum.add_paragraph()
    p1.text = "• 능동적인 자가 조절력 회복"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(14)
    p1.font.color.rgb = COLOR_WHITE
    p1.font.bold = True
    p1.space_after = Pt(6)

    p2 = tf_sum.add_paragraph()
    p2.text = "단순한 질병 부재를 넘어 스스로 몸과 마음, 뇌를 다각도로 보듬고 조화롭게 관리하는 자가 조절 훈련법(BOS)을 현장에 정착시키겠습니다."
    p2.font.name = FONT_NAME
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.space_after = Pt(16)

    p3 = tf_sum.add_paragraph()
    p3.text = "• 1대1 웰니스 피드백 로드맵 안내"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_WHITE
    p3.font.bold = True
    p3.space_after = Pt(6)

    p4 = tf_sum.add_paragraph()
    p4.text = "오늘 측정받으신 결과에 따라 개인별 최적화된 호흡, 명상, 기공, 운동 훈련 맞춤 로드맵을 배부하며 가이드해 드리겠습니다. 경청해 주셔서 감사합니다."
    p4.font.name = FONT_NAME
    p4.font.size = Pt(12)
    p4.font.color.rgb = COLOR_TEXT_MUTED
    p4.space_after = Pt(36)

    p5 = tf_sum.add_paragraph()
    p5.text = "💡 Q&A: 궁금하신 점을 편안하게 질문해 주십시오"
    p5.font.name = FONT_NAME
    p5.font.size = Pt(15)
    p5.font.color.rgb = COLOR_ACCENT
    p5.font.bold = True

    # PPT 파일 저장
    output_ppt = "교사_대상_3바디_7코드_AI_웰니스_강의안.pptx"
    try:
        prs.save(output_ppt)
        print(f"[완료] 교사 강의안 PPTX 생성 성공: {output_ppt}")
    except PermissionError:
        print(f"[알림] {output_ppt} 파일이 현재 사용 중이어서 덮어쓰기를 건너뜁니다.")

if __name__ == '__main__':
    generate_lecture_presentation()
