# 대한민국 건강 보고서 V2 파워포인트 자동 생성 스크립트
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 테마 색상 (프리미엄 다크 팔레트)
COLOR_DARK_BG = RGBColor(15, 23, 42)
COLOR_LIGHT_BG = RGBColor(248, 250, 252)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PRIMARY = RGBColor(79, 70, 229)
COLOR_SECONDARY = RGBColor(59, 130, 246)
COLOR_TEXT_DARK = RGBColor(30, 41, 59)
COLOR_TEXT_MUTED = RGBColor(100, 116, 139)
COLOR_ACCENT = RGBColor(245, 158, 11)
COLOR_SUCCESS = RGBColor(16, 185, 129)
COLOR_DANGER = RGBColor(239, 68, 68)

FONT_NAME = "Malgun Gothic"

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title, category):
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.3))
    p = cat_box.text_frame.paragraphs[0]
    p.text = category.upper()
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(10), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK

def card(slide, left, top, w, h, bg=COLOR_WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = bg
    s.line.color.rgb = RGBColor(226, 232, 240)
    s.line.width = Pt(1)
    return s

def fmt_cell(cell, text, size=11, bold=False, color=COLOR_TEXT_DARK, align=PP_ALIGN.CENTER):
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT_NAME
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

def style_header(table, cols):
    for i in range(cols):
        c = table.cell(0, i)
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_PRIMARY
        p = c.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

def main():
    with open('health_report_v2_data.json', 'r', encoding='utf-8') as f:
        data = json.load(f)

    summary = data['summary']
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ======= SLIDE 1: 표지 =======
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, COLOR_DARK_BG)

    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # V2 배지
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.5), Inches(1.8), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_ACCENT
    badge.line.fill.background()
    badge.text_frame.text = "제2판"
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = FONT_NAME
    p_badge.font.size = Pt(18)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_WHITE

    tb = s1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "CODEMAP AI 대국민 3Body 통합 건강 통계 보고서"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(16)
    p0.font.color.rgb = COLOR_SECONDARY
    p0.font.bold = True
    p0.space_after = Pt(12)

    p1 = tf.add_paragraph()
    p1.text = "대한민국 건강 보고서 V2"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = f"분석 대상: 고유 사용자 {summary['total']:,}명 (V1 대비 +{summary['total']-1955:,}명, 67.1% 증가)"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.space_after = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "기준 일자: 2026. 06. 17.  |  데이터: members_v4 (AI분석 완료 기준)"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED

    # ======= SLIDE 2: 전국 현황 요약 =======
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, COLOR_LIGHT_BG)
    add_header(s2, "📊 전국 점검 현황 요약", "OVERALL STATISTICS V2")

    # 카드 1: 총 분석 인원
    card(s2, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.0))
    tb1 = s2.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(3.0), Inches(0.4))
    tb1.text_frame.text = "🏆 AI분석 완료 인원"
    p = tb1.text_frame.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    tb_num = s2.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(3.0), Inches(1.0))
    tf_num = tb_num.text_frame
    pn = tf_num.paragraphs[0]
    pn.text = f"{summary['total']:,}명"
    pn.font.name = FONT_NAME
    pn.font.size = Pt(52)
    pn.font.bold = True
    pn.font.color.rgb = COLOR_DARK_BG

    pn2 = tf_num.add_paragraph()
    pn2.text = f"남성 {summary['male']:,}명 / 여성 {summary['female']:,}명"
    pn2.font.name = FONT_NAME
    pn2.font.size = Pt(13)
    pn2.font.color.rgb = COLOR_TEXT_MUTED

    # V1 대비 성장
    tb_growth = s2.shapes.add_textbox(Inches(1.1), Inches(3.8), Inches(3.0), Inches(0.5))
    pg = tb_growth.text_frame.paragraphs[0]
    pg.text = f"▲ V1 대비 +{summary['total']-1955:,}명 (67.1% 증가)"
    pg.font.name = FONT_NAME
    pg.font.size = Pt(14)
    pg.font.bold = True
    pg.font.color.rgb = COLOR_SUCCESS

    # 세부 정보
    tb_det = s2.shapes.add_textbox(Inches(1.1), Inches(4.6), Inches(3.0), Inches(1.5))
    tf_det = tb_det.text_frame
    tf_det.word_wrap = True
    for label, val in [
        ("PC 버전", f"{summary['pc_version']:,}명"),
        ("LITE 버전", f"{summary['lite_version']:,}명"),
        ("기타", f"{summary['other_version']:,}명"),
    ]:
        p = tf_det.add_paragraph()
        p.text = f"• {label}: {val}"
        p.font.name = FONT_NAME
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(4)

    # 카드 2: 3Body 핵심 점수
    card(s2, Inches(4.8), Inches(1.5), Inches(3.6), Inches(5.0))
    tb2 = s2.shapes.add_textbox(Inches(5.1), Inches(1.8), Inches(3.0), Inches(0.4))
    tb2.text_frame.text = "🎯 3Body 전체 평균 점수"
    p = tb2.text_frame.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY

    scores = [
        ("신체 (Body)", "70.5점", COLOR_SUCCESS),
        ("뇌 (Brain)", "77.3점", COLOR_SECONDARY),
        ("마음 (Mind)", "58.6점", COLOR_DANGER),
    ]
    y_pos = 2.5
    for label, val, color in scores:
        tb_s = s2.shapes.add_textbox(Inches(5.1), Inches(y_pos), Inches(3.0), Inches(0.8))
        tf_s = tb_s.text_frame
        ps1 = tf_s.paragraphs[0]
        ps1.text = label
        ps1.font.name = FONT_NAME
        ps1.font.size = Pt(12)
        ps1.font.color.rgb = COLOR_TEXT_MUTED

        ps2 = tf_s.add_paragraph()
        ps2.text = val
        ps2.font.name = FONT_NAME
        ps2.font.size = Pt(30)
        ps2.font.bold = True
        ps2.font.color.rgb = color
        y_pos += 1.2

    # 카드 3: 핵심 발견
    card(s2, Inches(8.8), Inches(1.5), Inches(3.6), Inches(5.0))
    tb3 = s2.shapes.add_textbox(Inches(9.1), Inches(1.8), Inches(3.0), Inches(0.4))
    tb3.text_frame.text = "🔑 핵심 발견"
    p = tb3.text_frame.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    findings = [
        "41.1% 이완 장애\n(1코드 최취약)",
        "20.4% 거북목 위험\n(V1 대비 +5.0%p)",
        "마음나이 > 실제나이\n(전 연령대 +10세)",
        "99.4% 에너지 순환\n관리 필요",
    ]
    y_pos = 2.5
    for finding in findings:
        tb_f = s2.shapes.add_textbox(Inches(9.1), Inches(y_pos), Inches(3.0), Inches(0.8))
        pf = tb_f.text_frame.paragraphs[0]
        pf.text = finding
        pf.font.name = FONT_NAME
        pf.font.size = Pt(12)
        pf.font.bold = True
        pf.font.color.rgb = COLOR_TEXT_DARK
        y_pos += 1.1

    # ======= SLIDE 3: 3Body 연령대별 =======
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, COLOR_LIGHT_BG)
    add_header(s3, "📈 3Body 연령대별 건강 지표", "AGE-BASED 3BODY ANALYSIS")

    age_data = data['age_threebody']
    ages = ['10대이하', '20대', '30대', '40대', '50대', '60대', '70대', '80대이상']
    rows = len(ages) + 1
    t_shape = s3.shapes.add_table(rows, 5, Inches(0.8), Inches(1.5), Inches(7.0), Inches(5.2))
    table = t_shape.table

    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(1.0)
    table.columns[2].width = Inches(1.4)
    table.columns[3].width = Inches(1.4)
    table.columns[4].width = Inches(1.4)

    headers = ["연령대", "인원", "신체(Body)", "뇌(Brain)", "마음(Mind)"]
    for i, h in enumerate(headers):
        fmt_cell(table.cell(0, i), h, 11, True, COLOR_WHITE)
    style_header(table, 5)

    for row_idx, ag in enumerate(ages, 1):
        d = age_data.get(ag, {})
        fmt_cell(table.cell(row_idx, 0), ag, bold=True, align=PP_ALIGN.LEFT)
        fmt_cell(table.cell(row_idx, 1), f"{d.get('count', 0):,}명")
        fmt_cell(table.cell(row_idx, 2), f"{d.get('body', '-')}점")
        fmt_cell(table.cell(row_idx, 3), f"{d.get('brain', '-')}점")
        # 마음 점수 강조 (60점 미만 빨간색)
        mind = d.get('mind', 0) or 0
        mind_color = COLOR_DANGER if mind < 60 else COLOR_TEXT_DARK
        fmt_cell(table.cell(row_idx, 4), f"{d.get('mind', '-')}점", color=mind_color, bold=True)

    # 인사이트 카드
    card(s3, Inches(8.2), Inches(1.5), Inches(4.3), Inches(5.2))
    tb_insight = s3.shapes.add_textbox(Inches(8.5), Inches(1.8), Inches(3.7), Inches(4.5))
    tf_i = tb_insight.text_frame
    tf_i.word_wrap = True

    insights = [
        ("💡 핵심 인사이트", COLOR_PRIMARY, True, 14),
        ("", COLOR_TEXT_DARK, False, 8),
        ("신체: 20대(76.3)에서 정점 → 80대(61.7)까지 완만 감소", COLOR_TEXT_DARK, False, 11),
        ("", COLOR_TEXT_DARK, False, 6),
        ("뇌: 40대(80.0)에서 정점 → 80대(74.4)까지 잘 보존", COLOR_TEXT_DARK, False, 11),
        ("", COLOR_TEXT_DARK, False, 6),
        ("마음: 전 연령대 60점 미만으로 만성 저조 (최대 과제)", COLOR_DANGER, True, 11),
        ("", COLOR_TEXT_DARK, False, 6),
        ("▶ 마음 점수가 유일하게 전 연령대에서 60점을 넘지 못하며, 현대인의 정서적 스트레스가 가장 심각한 건강 과제임을 입증합니다.", COLOR_TEXT_MUTED, False, 10),
    ]
    for idx, (text, color, bold, size) in enumerate(insights):
        if idx == 0:
            p = tf_i.paragraphs[0]
        else:
            p = tf_i.add_paragraph()
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color

    # ======= SLIDE 4: 7Code 취약 코드 =======
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, COLOR_LIGHT_BG)
    add_header(s4, "⚡ 7Code 웰니스 에너지 취약점 분석", "7CODE VULNERABILITY ANALYSIS")

    weak = data['weak_codes']
    total_weak = sum(weak.values())
    weak_items = sorted(weak.items(), key=lambda x: x[1], reverse=True)

    code_names = {
        '1': '1코드 (힐링 라이프 - 이완 장애)',
        '2': '2코드 (스트레스 관리 - 정서 긴장)',
        '3': '3코드 (감정 균형 - 무기력)',
        '4': '4코드 (활력 충전 - 에너지 고갈)',
        '5': '5코드 (집중·직관 - 인지 피로)',
        '6': '6코드 (통찰·지혜 - 두뇌 저하)',
        '7': '7코드 (의식·영성 - 정체성 혼란)',
    }

    t4_shape = s4.shapes.add_table(8, 4, Inches(0.8), Inches(1.5), Inches(7.0), Inches(4.0))
    t4 = t4_shape.table
    t4.columns[0].width = Inches(0.7)
    t4.columns[1].width = Inches(3.3)
    t4.columns[2].width = Inches(1.3)
    t4.columns[3].width = Inches(1.3)

    for i, h in enumerate(["순위", "코드 (증상)", "인원", "비율"]):
        fmt_cell(t4.cell(0, i), h, 11, True, COLOR_WHITE)
    style_header(t4, 4)

    for row_idx, (code, count) in enumerate(weak_items, 1):
        pct = round(count / total_weak * 100, 1)
        is_top = row_idx <= 2
        fmt_cell(t4.cell(row_idx, 0), str(row_idx), bold=is_top)
        fmt_cell(t4.cell(row_idx, 1), code_names.get(str(code), f'{code}코드'), align=PP_ALIGN.LEFT, bold=is_top)
        fmt_cell(t4.cell(row_idx, 2), f"{count:,}명", bold=is_top, color=COLOR_PRIMARY if is_top else COLOR_TEXT_DARK)
        fmt_cell(t4.cell(row_idx, 3), f"{pct}%", bold=is_top, color=COLOR_DANGER if is_top else COLOR_TEXT_DARK)

    # 인사이트 카드
    card(s4, Inches(8.2), Inches(1.5), Inches(4.3), Inches(4.0))
    tb4 = s4.shapes.add_textbox(Inches(8.5), Inches(1.8), Inches(3.7), Inches(3.5))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    for idx, (text, color, bold, size) in enumerate([
        ("🔥 국민 10명 중 4명", COLOR_DANGER, True, 16),
        ("이완 장애 (1코드 최취약)", COLOR_TEXT_DARK, True, 14),
        ("", COLOR_TEXT_DARK, False, 8),
        ("V1 (38.4%) → V2 (41.1%)", COLOR_ACCENT, True, 13),
        ("▲ 2.7%p 악화", COLOR_DANGER, True, 13),
        ("", COLOR_TEXT_DARK, False, 8),
        ("표본이 67% 증가해도 동일한 패턴이 강화되어, 이완 능력 결핍이 대한민국 국민의 가장 보편적 취약점임이 재확인되었습니다.", COLOR_TEXT_MUTED, False, 10),
    ]):
        p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color

    # ======= SLIDE 5: 자세 정렬 분석 =======
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, COLOR_LIGHT_BG)
    add_header(s5, "🦴 신체 관절 정렬 및 자세 위험도 분석", "POSTURE ALIGNMENT ANALYSIS (N=3,008)")

    posture = data['posture_grades']
    main_metrics = [
        ('거북목 (FHP) 및 경추 정렬', posture.get('거북목 (FHP) 및 경추 정렬', {})),
        ('어깨 / 골반 좌우 대칭', posture.get('어깨 / 골반 좌우 대칭', {})),
        ('측면 척추 정렬 (흉추/요추)', posture.get('측면 척추 정렬 (흉추/요추)', {})),
        ('하체 기저면 (무릎/다리/발목)', posture.get('하체 기저면 (무릎/다리/발목)', {})),
        ('귀-어깨-고관절-무릎 수직선 이탈', posture.get('귀-어깨-고관절-무릎 수직선 이탈', {})),
    ]

    t5_shape = s5.shapes.add_table(6, 5, Inches(0.8), Inches(1.5), Inches(7.4), Inches(3.5))
    t5 = t5_shape.table
    t5.columns[0].width = Inches(3.0)
    t5.columns[1].width = Inches(1.1)
    t5.columns[2].width = Inches(1.1)
    t5.columns[3].width = Inches(1.1)
    t5.columns[4].width = Inches(1.1)

    for i, h in enumerate(["자세 지표", "Good", "Fair", "Poor (위험)", "표본"]):
        fmt_cell(t5.cell(0, i), h, 11, True, COLOR_WHITE)
    style_header(t5, 5)

    for row_idx, (name, vals) in enumerate(main_metrics, 1):
        fmt_cell(t5.cell(row_idx, 0), name, align=PP_ALIGN.LEFT, bold=True)
        fmt_cell(t5.cell(row_idx, 1), f"{vals.get('Good_pct', 0)}%", color=COLOR_SUCCESS)
        fmt_cell(t5.cell(row_idx, 2), f"{vals.get('Fair_pct', 0)}%", color=COLOR_ACCENT)
        poor_pct = vals.get('Poor_pct', 0)
        poor_color = COLOR_DANGER if poor_pct >= 15 else COLOR_TEXT_DARK
        fmt_cell(t5.cell(row_idx, 3), f"{poor_pct}%", bold=True, color=poor_color)
        fmt_cell(t5.cell(row_idx, 4), f"{vals.get('sample', 0):,}명")

    # V1 대비 변화 카드
    card(s5, Inches(8.6), Inches(1.5), Inches(3.9), Inches(3.5))
    tb5 = s5.shapes.add_textbox(Inches(8.9), Inches(1.8), Inches(3.3), Inches(3.0))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    for idx, (text, color, bold, size) in enumerate([
        ("⚠️ V1 대비 악화 추이", COLOR_DANGER, True, 14),
        ("", COLOR_TEXT_DARK, False, 6),
        ("거북목 Poor: 15.4% → 20.4%", COLOR_TEXT_DARK, True, 12),
        ("(▲ +5.0%p)", COLOR_DANGER, True, 11),
        ("", COLOR_TEXT_DARK, False, 6),
        ("수직선 이탈 Poor: 15.3% → 22.4%", COLOR_TEXT_DARK, True, 12),
        ("(▲ +7.1%p)", COLOR_DANGER, True, 11),
        ("", COLOR_TEXT_DARK, False, 6),
        ("국민 5명 중 1명 이상이 전신 수직 정렬 위험 상태", COLOR_TEXT_MUTED, False, 10),
    ]):
        p = tf5.paragraphs[0] if idx == 0 else tf5.add_paragraph()
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color

    # ======= SLIDE 6: 건강나이 분석 =======
    s6 = prs.slides.add_slide(blank)
    set_bg(s6, COLOR_LIGHT_BG)
    add_header(s6, "🧬 건강나이(웰니스 에이지) 연령대별 비교", "WELLNESS AGE ESTIMATION")

    aging = data['age_aging']
    ages_for_aging = ['20대', '30대', '40대', '50대', '60대', '70대']

    t6_shape = s6.shapes.add_table(7, 7, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5))
    t6 = t6_shape.table
    t6.columns[0].width = Inches(1.2)
    for c in range(1, 7):
        t6.columns[c].width = Inches(1.7)

    for i, h in enumerate(["연령대", "실제 나이", "신체나이", "뇌나이", "마음나이", "얼굴나이", "종합나이"]):
        fmt_cell(t6.cell(0, i), h, 11, True, COLOR_WHITE)
    style_header(t6, 7)

    for row_idx, ag in enumerate(ages_for_aging, 1):
        d = aging.get(ag, {})
        actual = d.get('actual_age_avg', '-')
        fmt_cell(t6.cell(row_idx, 0), ag, bold=True, align=PP_ALIGN.LEFT)
        fmt_cell(t6.cell(row_idx, 1), f"{actual}세" if actual != '-' else '-')
        fmt_cell(t6.cell(row_idx, 2), f"{d.get('physical_age_avg', '-')}세")

        # 뇌나이 (실제보다 젊으면 파란색)
        brain_age = d.get('brain_age_avg')
        if brain_age and actual and actual != '-':
            brain_color = COLOR_SECONDARY if brain_age < actual else COLOR_DANGER
        else:
            brain_color = COLOR_TEXT_DARK
        fmt_cell(t6.cell(row_idx, 3), f"{d.get('brain_age_avg', '-')}세", color=brain_color, bold=True)

        # 마음나이 (실제보다 높으면 빨간색)
        mind_age = d.get('mind_age_avg')
        if mind_age and actual and actual != '-':
            mind_color = COLOR_DANGER if mind_age > actual else COLOR_SUCCESS
        else:
            mind_color = COLOR_TEXT_DARK
        fmt_cell(t6.cell(row_idx, 4), f"{d.get('mind_age_avg', '-')}세", color=mind_color, bold=True)

        fmt_cell(t6.cell(row_idx, 5), f"{d.get('face_age_avg', '-')}세")
        fmt_cell(t6.cell(row_idx, 6), f"{d.get('comprehensive_age_avg', '-')}세")

    # 인사이트
    card(s6, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5))
    tb6 = s6.shapes.add_textbox(Inches(1.1), Inches(5.5), Inches(11.0), Inches(1.0))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "💡 뇌나이는 실제 나이보다 항상 10~15세 젊고(파란색), 마음나이는 항상 10세 이상 높습니다(빨간색). 현대인의 뇌는 젊지만 마음은 늙어가고 있습니다."
    p6.font.name = FONT_NAME
    p6.font.size = Pt(13)
    p6.font.bold = True
    p6.font.color.rgb = COLOR_PRIMARY

    # ======= SLIDE 7: V1 vs V2 비교 =======
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, COLOR_LIGHT_BG)
    add_header(s7, "📊 V1 vs V2 핵심 변화 비교", "REPORT COMPARISON: V1 → V2")

    comparisons = [
        ("표본 규모", "1,955명", f"{summary['total']:,}명", "+67.1%", True),
        ("신체(Body) 평균", "71.4점", "70.5점", "▼ 0.9점", False),
        ("뇌(Brain) 평균", "76.5점", "77.3점", "▲ 0.8점", True),
        ("마음(Mind) 평균", "59.5점", "58.6점", "▼ 0.9점", False),
        ("1코드 취약 비율", "38.4%", "41.1%", "▲ 2.7%p", False),
        ("거북목 Poor", "15.4%", "20.4%", "▲ 5.0%p", False),
        ("수직선 이탈 Poor", "15.3%", "22.4%", "▲ 7.1%p", False),
        ("에너지 관리 필요", "99.7%", "99.4%", "동일", True),
    ]

    t7_shape = s7.shapes.add_table(len(comparisons)+1, 4, Inches(1.5), Inches(1.5), Inches(10.0), Inches(4.5))
    t7 = t7_shape.table
    t7.columns[0].width = Inches(3.0)
    t7.columns[1].width = Inches(2.3)
    t7.columns[2].width = Inches(2.3)
    t7.columns[3].width = Inches(2.4)

    for i, h in enumerate(["지표", "V1 (1,955명)", "V2 (3,267명)", "변화"]):
        fmt_cell(t7.cell(0, i), h, 12, True, COLOR_WHITE)
    style_header(t7, 4)

    for row_idx, (label, v1, v2, change, is_positive) in enumerate(comparisons, 1):
        fmt_cell(t7.cell(row_idx, 0), label, bold=True, align=PP_ALIGN.LEFT)
        fmt_cell(t7.cell(row_idx, 1), v1)
        fmt_cell(t7.cell(row_idx, 2), v2, bold=True, color=COLOR_PRIMARY)
        change_color = COLOR_SUCCESS if is_positive else COLOR_DANGER
        fmt_cell(t7.cell(row_idx, 3), change, bold=True, color=change_color)

    # 저장
    output = os.path.join(r'd:\antigravity_vibecoding\BT 3바디 ai테스트', '대한민국_건강보고서_V2.pptx')
    prs.save(output)
    print(f"✅ PPT 저장 완료: {output}")

if __name__ == '__main__':
    main()
