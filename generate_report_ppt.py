# 최종 지역별/지점별 점검 데이터를 바탕으로 파워포인트 리포트를 자동 생성하는 스크립트
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 테마 색상 정의 (Curated Premium Palette)
COLOR_DARK_BG = RGBColor(15, 23, 42)      # 딥 네이비 (표지 배경)
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # 슬레이트 라이트 그레이 (일반 슬라이드 배경)
COLOR_WHITE = RGBColor(255, 255, 255)     # 화이트 (카드/표 배경)
COLOR_PRIMARY = RGBColor(79, 70, 229)     # 인디고 퍼플 (메인 강조)
COLOR_SECONDARY = RGBColor(59, 130, 246)  # 블루 (서브 강조)
COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # 다크 그레이 (주 텍스트)
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 뮤티드 그레이 (부 텍스트)
COLOR_ACCENT = RGBColor(245, 158, 11)     # 오렌지/앰버 (포인트 강조)

FONT_NAME = "Malgun Gothic"

def add_header(slide, title_text, category_text="3Body 7Code AI분석 최종 결과"):
    """모든 상세 슬라이드에 일관된 헤더 영역을 만듭니다."""
    # 카테고리 태그 (상단 미니 텍스트)
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.0), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_NAME
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(10.0), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_DARK

def set_slide_background(slide, color):
    """슬라이드의 배경색을 설정합니다."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_card_shape(slide, left, top, width, height, bg_color=COLOR_WHITE):
    """깔끔한 모던 대시보드 카드 느낌의 직사각형 도형을 만듭니다."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    # 테두리 투명에 가깝게 회색 처리
    shape.line.color.rgb = RGBColor(226, 232, 240)
    shape.line.width = Pt(1)
    return shape

def format_cell(cell, text, font_size=11, bold=False, color=COLOR_TEXT_DARK, align=PP_ALIGN.CENTER):
    """표 셀의 텍스트 포맷을 정밀하게 제어합니다."""
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color

def style_table_header(table, cols):
    """표의 헤더 스타일을 고급 네이비/퍼플로 지정합니다."""
    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        # 텍스트 화이트/볼드로 세팅
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

def main():
    json_path = r"C:\Users\bange\.gemini\antigravity-ide\brain\6027f963-db0d-4493-9efe-8fd0848ae581\scratch\db_analysis_fixed_result.json"
    if not os.path.exists(json_path):
        print(f"Error: JSON file not found at {json_path}")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    summary = data['summary']
    stats = data['stats']

    # Presentation 초기화 및 16:9 와이드 비율 설정
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 기본 빈 슬라이드 레이아웃 (6: Blank)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: 표지 (다크 프리미엄 테마)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_DARK_BG)

    # 표지 메인 장식 (인디고 그라데이션 박스 느낌의 띠 디자인)
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_PRIMARY
    accent_bar.line.fill.background()

    # 타이틀 박스
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p_sub = tf1.paragraphs[0]
    p_sub.text = "본사 회원 데이터 통계 리포트"
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_SECONDARY
    p_sub.font.bold = True
    p_sub.space_after = Pt(14)

    p_main = tf1.add_paragraph()
    p_main.text = "3Body 7Code AI분석 최종 결과 보고"
    p_main.font.name = FONT_NAME
    p_main.font.size = Pt(44)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE
    p_main.space_after = Pt(20)

    p_info = tf1.add_paragraph()
    p_info.text = f"기준 일자: 2026. 06. 15.  |  대상 데이터: members_v4 (AI분석 완료 기준)"
    p_info.font.name = FONT_NAME
    p_info.font.size = Pt(13)
    p_info.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 2: 전국 집계 요약 (대시보드 카드 테마)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_LIGHT_BG)
    add_header(slide2, "📊 전국 점검 현황 요약", "OVERALL STATISTICS")

    # 카드 1: 분석완료 총합 (좌측 배치)
    create_card_shape(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
    c1_title = slide2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.5))
    c1_title.text_frame.text = "🏆 AI분석 완료 인원"
    c1_title.text_frame.paragraphs[0].font.name = FONT_NAME
    c1_title.text_frame.paragraphs[0].font.size = Pt(16)
    c1_title.text_frame.paragraphs[0].font.bold = True
    c1_title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    # 총 인원 크고 굵게 표시
    c1_number = slide2.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(1.2))
    tf_num = c1_number.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = f"{summary['grandTotal']:,}명"
    p_num.font.name = FONT_NAME
    p_num.font.size = Pt(54)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_DARK_BG
    
    p_num_sub = tf_num.add_paragraph()
    p_num_sub.text = f"(신규 {summary['totalNew']:,}명 / 기존 {summary['totalExisting']:,}명)"
    p_num_sub.font.name = FONT_NAME
    p_num_sub.font.size = Pt(14)
    p_num_sub.font.bold = True
    p_num_sub.font.color.rgb = COLOR_TEXT_MUTED

    # 세부 정보
    c1_detail = slide2.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(5.0), Inches(1.8))
    tf_det = c1_detail.text_frame
    def add_bullet(tf, label, val):
        p = tf.add_paragraph()
        p.text = f"• {label}: {val}"
        p.font.name = FONT_NAME
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)

    add_bullet(tf_det, "신규 회원 비중", f"{summary['totalNew'] / summary['grandTotal'] * 100:.1f}%")
    add_bullet(tf_det, "기존 회원 비중", f"{summary['totalExisting'] / summary['grandTotal'] * 100:.1f}%")
    add_bullet(tf_det, "기타(미분류 유형)", f"{summary['totalOther']:,}명")
    add_bullet(tf_det, "총 데이터 원본(분석완료)", f"{summary['totalAnalyzed']:,}명")

    # 카드 2: 데이터 정제 통계 (우측 배치)
    create_card_shape(slide2, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
    c2_title = slide2.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.0), Inches(0.5))
    c2_title.text_frame.text = "🧹 데이터 클리닝 및 제외 내역"
    c2_title.text_frame.paragraphs[0].font.name = FONT_NAME
    c2_title.text_frame.paragraphs[0].font.size = Pt(16)
    c2_title.text_frame.paragraphs[0].font.bold = True
    c2_title.text_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY

    c2_detail = slide2.shapes.add_textbox(Inches(7.1), Inches(2.6), Inches(5.0), Inches(3.4))
    tf_det2 = c2_detail.text_frame
    tf_det2.word_wrap = True
    
    p_desc = tf_det2.paragraphs[0]
    p_desc.text = "기기 오집계 및 분석 미완료 데이터를 제외하는 보정 알고리즘을 적용한 데이터 클리닝 결과입니다."
    p_desc.font.name = FONT_NAME
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_desc.space_after = Pt(20)

    add_bullet(tf_det2, "분석 대기 제외 (pending-)", f"{summary['pendingCount']:,}건")
    add_bullet(tf_det2, "기기 활성화 이전 데이터 제외 (천천동 국한)", f"{summary['excludedPreActivationCount']:,}건")
    add_bullet(tf_det2, "본사 교육/테스트 데이터 제외", f"{summary['excludedHQCount']:,}건")
    add_bullet(tf_det2, "지점 합동코드(EVT_) 재배정 완료", f"{summary['reallocatedEventCount']:,}건")

    # =========================================================================
    # SLIDE 3: 전국 지역별 순위 현황
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_LIGHT_BG)
    add_header(slide3, "📍 전국 지역별 실적 순위", "REGIONAL RANKINGS")

    # 지역별 데이터 요약 계산
    region_ranking = []
    for r_name, r_branches in stats.items():
        r_new = 0
        r_exist = 0
        for b_name, b_count in r_branches.items():
            r_new += b_count['new']
            r_exist += b_count['existing']
        region_ranking.append({
            'name': r_name,
            'new': r_new,
            'exist': r_exist,
            'total': r_new + r_exist
        })
    
    # 합계(total) 기준으로 내림차순 정렬
    region_ranking.sort(key=lambda x: x['total'], reverse=True)

    # 표 위치 및 배치
    # 총 23개 지역을 2개의 표로 병렬 배치하여 한눈에 보이게 설계
    table1_rows = 13  # 헤더1 + 12개 지역
    table2_rows = 12  # 헤더1 + 11개 지역
    cols = 5

    # 첫 번째 지역 순위 표 (1~12위)
    t1_left = Inches(0.8)
    t1_top = Inches(1.5)
    t1_width = Inches(5.6)
    t1_height = Inches(5.2)

    table1_shape = slide3.shapes.add_table(table1_rows, cols, t1_left, t1_top, t1_width, t1_height)
    table1 = table1_shape.table

    # 두 번째 지역 순위 표 (13~23위)
    t2_left = Inches(6.9)
    t2_top = Inches(1.5)
    t2_width = Inches(5.6)
    t2_height = Inches(4.8)

    table2_shape = slide3.shapes.add_table(table2_rows, cols, t2_left, t2_top, t2_width, t2_height)
    table2 = table2_shape.table

    # 컬럼 너비 설정
    for table_obj in [table1, table2]:
        table_obj.columns[0].width = Inches(0.7)  # 순위
        table_obj.columns[1].width = Inches(1.7)  # 지역명
        table_obj.columns[2].width = Inches(1.1)  # 합계
        table_obj.columns[3].width = Inches(1.0)  # 신규
        table_obj.columns[4].width = Inches(1.1)  # 기존

    # 헤더 텍스트 설정
    headers = ["순위", "지역명", "합계", "신규", "기존"]
    for t_idx, table_obj in enumerate([table1, table2]):
        for col_idx, h_text in enumerate(headers):
            format_cell(table_obj.cell(0, col_idx), h_text, 11, True, COLOR_WHITE)
        style_table_header(table_obj, cols)

    # 데이터 입력
    for idx, item in enumerate(region_ranking):
        rank = idx + 1
        if rank <= 12:
            table_obj = table1
            row_idx = rank
        else:
            table_obj = table2
            row_idx = rank - 12
        
        format_cell(table_obj.cell(row_idx, 0), str(rank))
        format_cell(table_obj.cell(row_idx, 1), item['name'], bold=True, align=PP_ALIGN.LEFT)
        format_cell(table_obj.cell(row_idx, 2), f"{item['total']}명", bold=True, color=COLOR_PRIMARY)
        format_cell(table_obj.cell(row_idx, 3), f"{item['new']}명")
        format_cell(table_obj.cell(row_idx, 4), f"{item['exist']}명")

    # =========================================================================
    # SLIDE 4: 전국 지점별 Top 20 순위 현황
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_LIGHT_BG)
    add_header(slide4, "🏆 전국 지점별 Top 20 실적 순위", "BRANCH TOP 20")

    # 모든 지점 데이터 평탄화
    branch_ranking = []
    for r_name, r_branches in stats.items():
        for b_name, b_count in r_branches.items():
            b_new = b_count['new']
            b_exist = b_count['existing']
            branch_ranking.append({
                'name': b_name,
                'region': r_name,
                'new': b_new,
                'exist': b_exist,
                'total': b_new + b_exist
            })

    # 합계(total) 기준 내림차순 정렬
    branch_ranking.sort(key=lambda x: x['total'], reverse=True)
    top_20_branches = branch_ranking[:20]

    # Top 20도 가독성을 위해 2열(1~10위, 11~20위)로 분할 배치
    b1_left = Inches(0.8)
    b1_top = Inches(1.5)
    b1_width = Inches(5.6)
    b1_height = Inches(4.8)

    table_b1_shape = slide4.shapes.add_table(11, 6, b1_left, b1_top, b1_width, b1_height)
    table_b1 = table_b1_shape.table

    b2_left = Inches(6.9)
    b2_top = Inches(1.5)
    b2_width = Inches(5.6)
    b2_height = Inches(4.8)

    table_b2_shape = slide4.shapes.add_table(11, 6, b2_left, b2_top, b2_width, b2_height)
    table_b2 = table_b2_shape.table

    branch_headers = ["순위", "지점명", "소속 지역", "합계", "신규", "기존"]
    for table_obj in [table_b1, table_b2]:
        table_obj.columns[0].width = Inches(0.6)  # 순위
        table_obj.columns[1].width = Inches(1.4)  # 지점명
        table_obj.columns[2].width = Inches(1.3)  # 소속 지역
        table_obj.columns[3].width = Inches(0.8)  # 합계
        table_obj.columns[4].width = Inches(0.7)  # 신규
        table_obj.columns[5].width = Inches(0.8)  # 기존
        
        for col_idx, h_text in enumerate(branch_headers):
            format_cell(table_obj.cell(0, col_idx), h_text, 11, True, COLOR_WHITE)
        style_table_header(table_obj, 6)

    for idx, item in enumerate(top_20_branches):
        rank = idx + 1
        if rank <= 10:
            table_obj = table_b1
            row_idx = rank
        else:
            table_obj = table_b2
            row_idx = rank - 10

        format_cell(table_obj.cell(row_idx, 0), str(rank))
        format_cell(table_obj.cell(row_idx, 1), item['name'], bold=True, align=PP_ALIGN.LEFT)
        format_cell(table_obj.cell(row_idx, 2), item['region'], align=PP_ALIGN.LEFT)
        format_cell(table_obj.cell(row_idx, 3), f"{item['total']}명", bold=True, color=COLOR_PRIMARY)
        format_cell(table_obj.cell(row_idx, 4), f"{item['new']}명")
        format_cell(table_obj.cell(row_idx, 5), f"{item['exist']}명")

    # =========================================================================
    # SLIDE 5~8: 지역별/지점별 상세 현황 (2열 카드 리스트 레이아웃)
    # =========================================================================
    # 모든 지역을 돌아가면서 상세 정보를 텍스트/표로 나열
    # 23개 지역을 4장의 슬라이드에 나누어 수록 (한 슬라이드당 6개 지역 내외)
    sortedRegions = sorted(stats.keys())
    regions_chunked = [sortedRegions[i:i+6] for i in range(0, len(sortedRegions), 6)]

    for slide_idx, chunk in enumerate(regions_chunked):
        slide_det = prs.slides.add_slide(blank_layout)
        set_slide_background(slide_det, COLOR_LIGHT_BG)
        add_header(slide_det, f"🔍 지역별 지점 상세 현황 ({slide_idx+1}/{len(regions_chunked)})", "DETAILED BRANCH STATISTICS")

        # 6개 지역 카드를 2열 3행 레이아웃으로 배치
        # 좌열: 0, 1, 2  |  우열: 3, 4, 5
        card_positions = [
            (Inches(0.8), Inches(1.5), Inches(5.6), Inches(1.6)),  # 1행 좌
            (Inches(0.8), Inches(3.3), Inches(5.6), Inches(1.6)),  # 2행 좌
            (Inches(0.8), Inches(5.1), Inches(5.6), Inches(1.6)),  # 3행 좌
            (Inches(6.9), Inches(1.5), Inches(5.6), Inches(1.6)),  # 1행 우
            (Inches(6.9), Inches(3.3), Inches(5.6), Inches(1.6)),  # 2행 우
            (Inches(6.9), Inches(5.1), Inches(5.6), Inches(1.6)),  # 3행 우
        ]

        for item_idx, r_name in enumerate(chunk):
            left, top, width, height = card_positions[item_idx]
            
            # 카드 배경 상자
            create_card_shape(slide_det, left, top, width, height)
            
            # 카드 타이틀 (지역명)
            title_text = f"📍 {r_name}"
            # 지역 합계 계산
            r_branches = stats[r_name]
            r_new = sum(b['new'] for b in r_branches.values())
            r_exist = sum(b['existing'] for b in r_branches.values())
            r_total = r_new + r_exist
            
            title_box = slide_det.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4))
            tf_ct = title_box.text_frame
            p_ct = tf_ct.paragraphs[0]
            p_ct.text = f"{title_text}  (합계 {r_total}명 - 신규 {r_new} / 기존 {r_exist})"
            p_ct.font.name = FONT_NAME
            p_ct.font.size = Pt(13)
            p_ct.font.bold = True
            p_ct.font.color.rgb = COLOR_PRIMARY

            # 카드 바디 (지점들 텍스트 나열)
            # 지점 수가 많으면 가독성을 위해 간략하게 텍스트로 축약
            body_box = slide_det.shapes.add_textbox(left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(1.1))
            tf_cb = body_box.text_frame
            tf_cb.word_wrap = True
            
            sorted_b = sorted(r_branches.items(), key=lambda x: x[1]['new'] + x[1]['existing'], reverse=True)
            
            # 지점 문자열 목록 빌드
            b_strings = []
            for b_name, b_cnt in sorted_b:
                b_sum = b_cnt['new'] + b_cnt['existing']
                b_strings.append(f"{b_name}({b_sum}명)")
            
            p_body = tf_cb.paragraphs[0]
            p_body.text = "   ".join(b_strings[:10])  # 최대 10개 지점까지 한 줄에 노출
            p_body.font.name = FONT_NAME
            p_body.font.size = Pt(10)
            p_body.font.color.rgb = COLOR_TEXT_DARK
            p_body.space_before = Pt(4)

            # 지점이 10개 이상이면 더보기 표시
            if len(b_strings) > 10:
                p_more = tf_cb.add_paragraph()
                p_more.text = f"...외 {len(b_strings)-10}개 지점"
                p_more.font.name = FONT_NAME
                p_more.font.size = Pt(9)
                p_more.font.color.rgb = COLOR_TEXT_MUTED
                p_more.font.italic = True

    # 프레젠테이션 저장
    output_filename = "3body_ai_analysis_report.pptx"
    output_path = os.path.join("d:\\antigravity_vibecoding\\BT 3바디 ai테스트", output_filename)
    prs.save(output_path)
    print(f"Presentation saved successfully at {output_path}")

if __name__ == "__main__":
    main()
