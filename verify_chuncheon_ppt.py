# 춘천교육청 PPT 보고서 구조 검증 스크립트
from pptx import Presentation

prs = Presentation("춘천교육청_코드맵AI_결과보고서.pptx")
print(f"총 슬라이드 수: {len(prs.slides)}")
print(f"슬라이드 크기: {prs.slide_width} x {prs.slide_height}")

for i, slide in enumerate(prs.slides, 1):
    shapes_count = len(slide.shapes)
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t and len(t) > 2:
                    texts.append(t[:60])
    print(f"\n  Slide {i}: {shapes_count}개 도형")
    for t in texts[:5]:
        print(f"    - {t}")

print("\n검증 완료.")
