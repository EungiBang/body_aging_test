# BTC 단체 건강 점검 결과 PPT 보고서 자동 생성 스크립트
import argparse
import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

C_DARK = RGBColor(15, 23, 42); C_LIGHT = RGBColor(248, 250, 252); C_WHITE = RGBColor(255, 255, 255)
C_PRIMARY = RGBColor(31, 73, 125); C_SECONDARY = RGBColor(59, 130, 246)
C_TEXT = RGBColor(30, 41, 59); C_MUTED = RGBColor(100, 116, 139)
C_ACCENT = RGBColor(245, 158, 11); C_SUCCESS = RGBColor(16, 185, 129); C_DANGER = RGBColor(239, 68, 68)
C_BORDER = RGBColor(226, 232, 240); C_INDIGO = RGBColor(79, 70, 229)
C_FUCHSIA = RGBColor(192, 38, 211); C_SLATE = RGBColor(148, 163, 184)
FONT = "Malgun Gothic"

def bg(s, c): f = s.background.fill; f.solid(); f.fore_color.rgb = c
def box(s, l, t, w, h, text="", sz=11, b=False, c=C_TEXT, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.text = text; p.font.name = FONT; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c; p.alignment = a
    return tb
def card(s, l, t, w, h, bg_c=C_WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h); sh.fill.solid(); sh.fill.fore_color.rgb = bg_c; sh.line.color.rgb = C_BORDER; sh.line.width = Pt(1)
    return sh
def header(s, title, sub):
    box(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3), sub.upper(), 9.5, True, C_SECONDARY)
    box(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.6), title, 22, True, C_TEXT)
def fmt(cell, text, sz=10, b=False, c=C_TEXT, a=PP_ALIGN.CENTER):
    cell.text_frame.text = ""; p = cell.text_frame.paragraphs[0]
    p.text = str(text) if text is not None else "-"; p.alignment = a; p.font.name = FONT; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c
def shdr(table, cols):
    for i in range(cols):
        c = table.cell(0, i); c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
        p = c.text_frame.paragraphs[0]; p.font.color.rgb = C_WHITE; p.font.bold = True
def kpi(s, l, t, w, h, label, value, sub="", bg_c=C_WHITE, vc=C_TEXT):
    card(s, l, t, w, h, bg_c)
    box(s, l+Inches(0.12), t+Inches(0.08), w-Inches(0.24), Inches(0.2), label, 8.5, True, C_MUTED, PP_ALIGN.CENTER)
    box(s, l+Inches(0.12), t+Inches(0.28), w-Inches(0.24), Inches(0.35), str(value), 18, True, vc, PP_ALIGN.CENTER)
    if sub: box(s, l+Inches(0.12), t+Inches(0.62), w-Inches(0.24), Inches(0.18), sub, 7.5, False, C_MUTED, PP_ALIGN.CENTER)
def diff_str(v, n):
    if v is None or n is None: return "-"
    d = round(v - n, 1); s = "+" if d > 0 else ""
    return f"전국 {n}점 ({s}{d})"

def generate(input_file, output_file):
    with open(input_file, "r", encoding="utf-8") as f:
        data = json.load(f)

    dm = data["demographics"]; tb = data["three_body"]; tba = data["three_body_by_age"]
    ha = data["health_age"]; sc = data["seven_code"]; wc = data["weakest_code_analysis"]
    pa = data["posture_analysis"]; na = data["national_avg"]
    total = dm["total"]; hao = ha["overall"]
    institution = data.get("institution", "단체")
    test_date = data.get("test_date", "")

    # 날짜 포매팅
    date_parts = test_date.split("-")
    date_display = f"{date_parts[0]}년 {int(date_parts[1])}월 {int(date_parts[2])}일" if len(date_parts) == 3 else test_date

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # === SLIDE 1: 표지 ===
    s = prs.slides.add_slide(blank); bg(s, C_DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_INDIGO; bar.line.fill.background()
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(2.6), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = C_SECONDARY; badge.line.fill.background()
    badge.text_frame.text = "CODEMAP AI  HEALTH REPORT"
    p = badge.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.name = FONT; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = C_WHITE
    box(s, Inches(1.2), Inches(2.7), Inches(11), Inches(1.0), f"{institution} 3Body 7Code 건강 점검 결과 보고서", 34, True, C_WHITE)
    box(s, Inches(1.2), Inches(3.8), Inches(11), Inches(0.5), "AI 기반 신체 . 뇌 . 마음 통합 건강 분석 리포트", 16, False, C_SECONDARY)
    for i, line in enumerate([f"점검 일자     {date_display}", f"점검 인원     {total}명 (남성 {dm['male_count']}명 / 여성 {dm['female_count']}명)", f"대상 기관     {institution}", "점검 방식     코드맵AI 3Body 7Code (LITE 버전)"]):
        box(s, Inches(1.2), Inches(4.8 + i * 0.4), Inches(11), Inches(0.35), line, 13, False, C_SLATE)
    box(s, Inches(1.2), Inches(6.6), Inches(11), Inches(0.4), "BTC 3Body Lab  |  CODEMAP AI  |  Confidential", 10, True, C_MUTED)

    # === SLIDE 2: 점검 개요 ===
    s = prs.slides.add_slide(blank); bg(s, C_LIGHT); header(s, "점검 개요 및 참여자 분석", "OVERVIEW")
    kpi(s, Inches(0.6), Inches(1.5), Inches(2.2), Inches(0.85), "총 점검 인원", f"{total}명", "", RGBColor(238,242,255), C_INDIGO)
    kpi(s, Inches(3.0), Inches(1.5), Inches(2.2), Inches(0.85), "남성", f"{dm['male_count']}명", f"{round(dm['male_count']/total*100,1)}%", C_WHITE, C_SECONDARY)
    kpi(s, Inches(5.4), Inches(1.5), Inches(2.2), Inches(0.85), "여성", f"{dm['female_count']}명", f"{round(dm['female_count']/total*100,1)}%", C_WHITE, C_FUCHSIA)
    kpi(s, Inches(7.8), Inches(1.5), Inches(2.2), Inches(0.85), "평균 나이", f"{dm['avg_age']}세", "", C_WHITE, C_TEXT)
    kpi(s, Inches(10.2), Inches(1.5), Inches(2.6), Inches(0.85), "점검일", test_date.replace("-", "."), institution, C_WHITE, C_ACCENT)
    age_order = ["20대","30대","40대","50대","60대","70대"]
    age_data = [(ag, dm["age_groups"].get(ag,0)) for ag in age_order if dm["age_groups"].get(ag,0) > 0]
    if age_data:
        tbl = s.shapes.add_table(3, len(age_data), Inches(0.6), Inches(2.8), Inches(12.2), Inches(1.0)).table
        for i, (ag, cnt) in enumerate(age_data):
            fmt(tbl.cell(0,i),ag,10,True,C_WHITE); fmt(tbl.cell(1,i),f"{cnt}명",14,True,C_INDIGO); fmt(tbl.cell(2,i),f"{round(cnt/total*100,1)}%",9,False,C_MUTED)
        shdr(tbl, len(age_data))
    card(s, Inches(0.6), Inches(4.2), Inches(12.2), Inches(3.0))
    box(s, Inches(0.9), Inches(4.4), Inches(11.6), Inches(2.7),
        f"코드맵AI 3Body 7Code 점검 시스템은 AI 기반으로 신체(Body), 뇌(Brain), 마음(Mind)의 3가지 영역을 실시간으로 측정하여 개인별 건강 프로파일을 생성하는 통합 건강관리 솔루션입니다.\n\n"
        f"- 신체(Body) 분석: TensorFlow.js 기반 실시간 자세 스캔 + AI 체형 정렬 분석\n"
        f"- 뇌(Brain) 분석: 시각 반응 속도, 패턴 인지, 기억력 테스트 기반 인지 기능 측정\n"
        f"- 마음(Mind) 분석: 설문 + 생체 반응 융합 분석을 통한 정서적 건강 상태 추정\n"
        f"- 7Code 에너지: 7가지 에너지 코드로 전신 웰니스 상태를 코드화\n\n"
        f"본 보고서는 {institution} 소속 {total}명의 전체 데이터를 Firebase 서버에서 직접 추출하여 분석하고, 전국 {na['total_members']:,}명의 평균 데이터와 비교합니다.", 10.5, False, C_TEXT)

    # === SLIDE 3: 3Body ===
    s = prs.slides.add_slide(blank); bg(s, C_LIGHT); header(s, "3Body 통합 건강 점수 분석", "3BODY ANALYSIS")
    tbo = tb["overall"]
    for i, (label, key, color) in enumerate([("신체 (Body)","body",C_INDIGO),("뇌 (Brain)","brain",C_ACCENT),("마음 (Mind)","mind",C_FUCHSIA)]):
        v = tbo[key]; n = na["three_body"][key]
        kpi(s, Inches(0.6+i*4.2), Inches(1.5), Inches(3.8), Inches(1.1), label, f"{v}점", diff_str(v, n), C_WHITE, color)
    tbl = s.shapes.add_table(4, 4, Inches(0.6), Inches(3.0), Inches(6.0), Inches(1.8)).table
    for i, h in enumerate(["구분","신체 (Body)","뇌 (Brain)","마음 (Mind)"]): fmt(tbl.cell(0,i),h,10,True,C_WHITE)
    shdr(tbl, 4)
    for r, (label, vals) in enumerate([(f"{institution}", tb["overall"]),("남성",tb["male"]),("여성",tb["female"])],1):
        fmt(tbl.cell(r,0),label,10,True,C_TEXT)
        for ci, k in enumerate(["body","brain","mind"],1): fmt(tbl.cell(r,ci),f"{vals[k]}점" if vals.get(k) else "-",10,False,C_TEXT)
    tbl2 = s.shapes.add_table(3, 4, Inches(7.0), Inches(3.0), Inches(5.8), Inches(1.3)).table
    for i, h in enumerate(["비교","신체","뇌","마음"]): fmt(tbl2.cell(0,i),h,10,True,C_WHITE)
    shdr(tbl2, 4)
    fmt(tbl2.cell(1,0),institution,10,True,C_TEXT); fmt(tbl2.cell(2,0),"전국 평균",10,True,C_MUTED)
    for ci, k in enumerate(["body","brain","mind"],1):
        v=tbo[k]; n=na["three_body"][k]; d=round(v-n,1); color = C_SUCCESS if d>0 else C_DANGER
        fmt(tbl2.cell(1,ci),f"{v}점",11,True,color); fmt(tbl2.cell(2,ci),f"{n}점",10,False,C_MUTED)
    if tba:
        rows_a = len(tba)+1
        tbl3 = s.shapes.add_table(rows_a, 5, Inches(0.6), Inches(5.2), Inches(12.2), Inches(0.3+rows_a*0.32)).table
        for i, h in enumerate(["연령대","인원","신체","뇌","마음"]): fmt(tbl3.cell(0,i),h,10,True,C_WHITE)
        shdr(tbl3, 5)
        for r, (ag, v) in enumerate(tba.items(), 1):
            fmt(tbl3.cell(r,0),ag,10,True,C_TEXT); fmt(tbl3.cell(r,1),f"{v['count']}명",10,False,C_TEXT)
            fmt(tbl3.cell(r,2),f"{v['body']}점",10,False,C_TEXT); fmt(tbl3.cell(r,3),f"{v['brain']}점",10,False,C_TEXT); fmt(tbl3.cell(r,4),f"{v['mind']}점",10,False,C_TEXT)

    # === SLIDE 4: 건강나이 ===
    s = prs.slides.add_slide(blank); bg(s, C_LIGHT); header(s, "건강나이 (웰니스 에이지) 분석", "WELLNESS AGE")
    for i, (label, val, unit, color) in enumerate([("실제 나이",hao["avg_real_age"],"세",C_TEXT),("신체 나이",hao["avg_physical_age"],"세",C_DANGER),("뇌 나이",hao["avg_brain_age"],"세",C_SUCCESS),("마음 나이",hao["avg_mind_age"],"세",C_FUCHSIA),("종합 나이",hao["avg_comprehensive_age"],"세",C_INDIGO),("종합 점수",hao["avg_overall_score"],"점",C_ACCENT)]):
        sub = ""
        if label in ["신체 나이","뇌 나이","마음 나이"] and val and hao["avg_real_age"]:
            gap = round(val-hao["avg_real_age"],1); sub = f"실제 대비 {'+' if gap>0 else ''}{gap}세"
        kpi(s, Inches(0.6+(i%6)*2.1), Inches(1.5), Inches(1.9), Inches(0.85), label, f"{val}{unit}" if val else "-", sub, C_WHITE, color)
    tbl = s.shapes.add_table(3, 7, Inches(0.6), Inches(2.7), Inches(12.2), Inches(1.2)).table
    for i, h in enumerate(["구분","실제 나이","신체 나이","뇌 나이","마음 나이","종합 나이","종합 점수"]): fmt(tbl.cell(0,i),h,9,True,C_WHITE)
    shdr(tbl, 7)
    for r, (label, d) in enumerate([("남성",ha["male"]),("여성",ha["female"])],1):
        fmt(tbl.cell(r,0),label,10,True,C_TEXT)
        for ci, (fld, u) in enumerate(zip(["avg_real_age","avg_physical_age","avg_brain_age","avg_mind_age","avg_comprehensive_age","avg_overall_score"],["세","세","세","세","세","점"]),1):
            v = d.get(fld); fmt(tbl.cell(r,ci),f"{v}{u}" if v else "-",10,False,C_TEXT)
    card(s, Inches(0.6), Inches(4.2), Inches(12.2), Inches(1.2))
    mind_gap = round(hao["avg_mind_age"]-hao["avg_real_age"],1) if hao.get("avg_mind_age") and hao.get("avg_real_age") else 0
    brain_gap = round(hao["avg_brain_age"]-hao["avg_real_age"],1) if hao.get("avg_brain_age") and hao.get("avg_real_age") else 0
    phys_gap = round(hao["avg_physical_age"]-hao["avg_real_age"],1) if hao.get("avg_physical_age") and hao.get("avg_real_age") else 0
    box(s, Inches(0.9), Inches(4.35), Inches(11.6), Inches(1.0),
        f"핵심 인사이트:  실제 나이 {hao['avg_real_age']}세 대비, 신체 나이 {'+'if phys_gap>0 else ''}{phys_gap}세, "
        f"뇌 나이 {'+'if brain_gap>0 else ''}{brain_gap}세, 마음 나이 {'+'if mind_gap>0 else ''}{mind_gap}세. "
        f"{'마음 나이가 가장 빠르게 노화되고 있어 스트레스 관리와 이완 훈련이 시급합니다.' if mind_gap > max(phys_gap, brain_gap) else '신체 나이가 가장 빠르게 노화되고 있어 체형 교정 및 운동 습관 개선이 필요합니다.'}", 11, False, C_TEXT)

    # === SLIDE 5: 7Code ===
    s = prs.slides.add_slide(blank); bg(s, C_LIGHT); header(s, "7Code 에너지 웰니스 분석", "7CODE ENERGY ANALYSIS")
    code_labels = ["C1 기초에너지","C2 감정흐름","C3 추진력","C4 정서안정","C5 소통","C6 집중통찰","C7 삶의방향"]
    tbl = s.shapes.add_table(4, 8, Inches(0.6), Inches(1.5), Inches(12.2), Inches(1.8)).table
    fmt(tbl.cell(0,0),"구분",10,True,C_WHITE)
    for i, cl in enumerate(code_labels): fmt(tbl.cell(0,i+1),cl,8.5,True,C_WHITE)
    shdr(tbl, 8)
    for r, (rl, rd) in enumerate([(institution,sc["overall"]),("전국 평균",na["seven_code"])],1):
        fmt(tbl.cell(r,0),rl,10,True,C_TEXT)
        for i in range(7): fmt(tbl.cell(r,i+1),f"{rd[f'c{i+1}']}점",10,False,C_TEXT)
    fmt(tbl.cell(3,0),"차이",10,True,C_TEXT)
    for i in range(7):
        v=sc["overall"][f"c{i+1}"]; n=na["seven_code"][f"c{i+1}"]; d=round(v-n,1); color = C_SUCCESS if d>0 else C_DANGER
        fmt(tbl.cell(3,i+1),f"{'+'if d>0 else ''}{d}",10,True,color)
    box(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.35), "최취약 코드 TOP 3 (개인별 최하위 점수 코드로 분류)", 12, True, C_TEXT)
    sorted_wc = sorted(wc.items(), key=lambda x: x[1]["count"], reverse=True)
    for i, (cn, info) in enumerate(sorted_wc[:3]):
        x = Inches(0.6+i*4.2); rc = [C_DANGER,C_ACCENT,C_SECONDARY][i]
        card(s, x, Inches(4.1), Inches(3.8), Inches(1.3))
        box(s, x+Inches(0.15), Inches(4.15), Inches(0.4), Inches(0.35), f"#{i+1}", 13, True, rc, PP_ALIGN.CENTER)
        box(s, x+Inches(0.6), Inches(4.15), Inches(3.0), Inches(0.3), info["label"], 10, True, C_TEXT)
        box(s, x+Inches(0.6), Inches(4.5), Inches(3.0), Inches(0.3), f"{info['count']}명 ({info['ratio']}%)", 16, True, rc, PP_ALIGN.CENTER)
        nr = na["weakest_code_ratio"].get(int(cn), 0)
        box(s, x+Inches(0.6), Inches(4.9), Inches(3.0), Inches(0.25), f"전국 평균: {nr}%", 9, False, C_MUTED, PP_ALIGN.CENTER)
    card(s, Inches(0.6), Inches(5.7), Inches(12.2), Inches(1.5))
    if sorted_wc:
        tw = sorted_wc[0]
        box(s, Inches(0.9), Inches(5.85), Inches(11.6), Inches(1.2),
            f"핵심 인사이트:  {institution}의 {tw[1]['ratio']}%가 코드{tw[0]}({tw[1]['label']}) 영역에서 가장 취약합니다. "
            f"이는 전국 평균({na['weakest_code_ratio'].get(int(tw[0]),0)}%) 대비 높은 수치이며, 해당 영역의 집중 관리가 필요합니다.", 10.5, False, C_TEXT)

    # === SLIDE 6: 자세 ===
    if pa:
        s = prs.slides.add_slide(blank); bg(s, C_LIGHT); header(s, "신체 자세 정렬 상태 분석", "POSTURE ANALYSIS")
        pk = list(pa.keys()); rows_p = len(pk)+1
        tbl = s.shapes.add_table(rows_p, 5, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.35+rows_p*0.45)).table
        for i, h in enumerate(["자세 측정 지표","Good (정상)","Fair (주의)","Poor (위험)","표본"]): fmt(tbl.cell(0,i),h,10,True,C_WHITE)
        shdr(tbl, 5)
        for r, nm in enumerate(pk, 1):
            st = pa[nm]; fmt(tbl.cell(r,0),nm,9,False,C_TEXT,PP_ALIGN.LEFT); fmt(tbl.cell(r,1),f"{st['Good']}%",10,False,C_SUCCESS)
            fmt(tbl.cell(r,2),f"{st['Fair']}%",10,False,C_ACCENT)
            pc = C_DANGER if st["Poor"]>10 else C_TEXT; fmt(tbl.cell(r,3),f"{st['Poor']}%",10,True,pc)
            fmt(tbl.cell(r,4),f"{st['total']}명",10,False,C_MUTED)
        nat_p = na["posture"]; compare = list(nat_p.keys()); c_cnt = len(compare)+1
        tbl2 = s.shapes.add_table(c_cnt, 4, Inches(0.6), Inches(1.5+rows_p*0.45+0.5), Inches(12.2), Inches(0.35+c_cnt*0.4)).table
        for i, h in enumerate(["지표",f"{institution} Poor","전국 Poor","비교"]): fmt(tbl2.cell(0,i),h,10,True,C_WHITE)
        shdr(tbl2, 4)
        for r, nn in enumerate(compare, 1):
            cc_val = None
            for ppk, ppv in pa.items():
                if nn[:4] in ppk: cc_val = ppv["Poor"]; break
            nv = nat_p[nn]["Poor"]
            fmt(tbl2.cell(r,0),nn,8,False,C_TEXT,PP_ALIGN.LEFT); fmt(tbl2.cell(r,1),f"{cc_val}%" if cc_val is not None else "-",10,True,C_TEXT)
            fmt(tbl2.cell(r,2),f"{nv}%",10,False,C_MUTED)
            if cc_val is not None:
                d = round(cc_val-nv,1); color = C_DANGER if d>0 else C_SUCCESS
                fmt(tbl2.cell(r,3),f"{'+'if d>0 else ''}{d}%p",10,True,color)

    # === SLIDE 7: 종합 ===
    s = prs.slides.add_slide(blank); bg(s, C_DARK)
    box(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.3), "CONCLUSION", 10, True, C_SECONDARY)
    box(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.6), "종합 인사이트 및 건강 관리 제언", 24, True, C_WHITE)

    findings = []
    # 동적 인사이트 생성
    ages = [("마음 나이", mind_gap, C_FUCHSIA), ("신체 나이", phys_gap, C_DANGER), ("뇌 나이", brain_gap, C_SUCCESS)]
    worst_age = max(ages, key=lambda x: x[1])
    best_age = min(ages, key=lambda x: x[1])
    findings.append((f"{worst_age[0]}가 실제보다 {abs(worst_age[1])}세 {'높습니다' if worst_age[1]>0 else '젊습니다'}",
        f"평균 실제 나이 {hao['avg_real_age']}세 대비 {worst_age[0]} {'+' if worst_age[1]>0 else ''}{worst_age[1]}세로, 해당 영역의 노화가 가장 빠릅니다.",worst_age[2]))
    if sorted_wc:
        tw = sorted_wc[0]
        findings.append((f"코드{tw[0]}({tw[1]['label']}) {tw[1]['ratio']}%가 최취약",
            f"{total}명 중 {tw[1]['count']}명이 해당 영역에서 최하위 점수를 기록했습니다. 전국 평균({na['weakest_code_ratio'].get(int(tw[0]),0)}%) 대비 높은 수치입니다.",C_DANGER))
    findings.append((f"{best_age[0]} 양호, 관련 역량 우수",
        f"{best_age[0]} {'+' if best_age[1]>0 else ''}{best_age[1]}세로, 해당 영역의 건강 상태가 양호합니다.",C_SUCCESS))
    findings.append(("정기 점검 및 조직 건강 프로그램 도입 제안",
        f"6개월 간격의 코드맵AI 정기 점검으로 건강 변화를 추적하고, 조직 차원의 이완 훈련, 점심 스트레칭, 뇌 활성화 훈련 도입을 제안합니다.",C_SECONDARY))

    for i, (title, desc, color) in enumerate(findings[:5]):
        y = Inches(1.5+i*1.25)
        c = card(s, Inches(0.6), y, Inches(12.2), Inches(1.1), RGBColor(30,41,59))
        c.line.color.rgb = RGBColor(51,65,85)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.15), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        box(s, Inches(1.0), y+Inches(0.1), Inches(11.5), Inches(0.3), title, 13, True, C_WHITE)
        box(s, Inches(1.0), y+Inches(0.45), Inches(11.5), Inches(0.55), desc, 10.5, False, C_SLATE)

    box(s, Inches(0.6), Inches(7.0), Inches(12.2), Inches(0.3),
        "본 보고서의 모든 데이터는 비의료 건강관리서비스 가이드라인을 준수하며, 의료행위를 대체하지 않습니다.", 8, False, C_MUTED, PP_ALIGN.CENTER)

    prs.save(output_file)
    print(f"Success! PPT written to: {output_file} ({len(prs.slides)} slides)")

def main():
    parser = argparse.ArgumentParser(description="BTC 단체 건강 점검 결과 PPT 보고서 생성")
    parser.add_argument("--input", required=True, help="분석 결과 JSON 파일 경로")
    parser.add_argument("--output", required=True, help="PPT 출력 파일 경로")
    args = parser.parse_args()
    generate(args.input, args.output)

if __name__ == "__main__":
    main()
